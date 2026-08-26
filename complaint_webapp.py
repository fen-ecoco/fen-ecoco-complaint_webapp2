import io
import json
import re
import zipfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional

import matplotlib.pyplot as plt
import pandas as pd
import plotly.express as px
import streamlit as st
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import gspread
    from google.oauth2.service_account import Credentials
except Exception:
    gspread = None
    Credentials = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None


st.set_page_config(page_title="ECOCO 客訴分析平台", page_icon="📊", layout="wide")

# ── 自動化核心（automation/ 套件；與無介面排程 automation/cli.py 共用同一套邏輯）──
from automation import config as auto_config
from automation.classifier import build_default as build_classifier
from automation.columns import detect_columns
from automation.core import (
    META_COLUMNS,
    AnalysisConfig,
    analyze_dataframe as core_analyze_dataframe,
    make_unique_columns,
    review_summary,
    visible_columns,
)
from automation.rules import _is_valid_pair, analyze_complaint
from automation.taxonomy import (
    DEPT_MAP,
    DEPT_OPTIONS,
    DETAIL_OPTIONS,
    TOPIC_DETAIL_MAP,
    TYPE_OPTIONS,
)
from automation.text import (
    lower_english,
    mask_phone_value,
    mask_sensitive_df,
    mask_sensitive_text,
    normalize_problem_labels,
    PHONE_COL_HINTS,
)



# ── ECOCO 品牌色（Pantone 對應）──────────────────────────────
BRAND_ORANGE  = "#FF5000"   # Pantone Orange 021 C  → 營運部
BRAND_BLUE    = "#060E9F"   # Pantone Blue 072 C    → 資訊部 / 主圖色
BRAND_YELLOW  = "#FFCE00"   # Pantone 116 C         → 行銷部
BRAND_LBLUE   = "#8EB9C9"   # Pantone 550 C
BRAND_BEIGE   = "#FAE0B8"   # Pantone P17-2 C
BRAND_TEAL    = "#0076A9"   # Pantone 7690 C
BRAND_WHITE   = "#FFFFFF"   # Pantone White C

# 部門固定色（Plotly color_discrete_map 用）
DEPT_COLOR_MAP: dict[str, str] = {
    "營運部": BRAND_ORANGE,
    "行銷部": BRAND_YELLOW,
    "資訊部": BRAND_BLUE,
    "研發部": BRAND_TEAL,
    "廠務部": BRAND_LBLUE,
    "人資部": BRAND_BEIGE,
    "企劃部": "#A0C878",
    "財務部": "#C8A0E0",
    "開發部": "#E0C8A0",
    "總經理室": "#A0E0C8",
    "未分配":  "#CCCCCC",
    "":        "#CCCCCC",
}

# 圓餅圖 / 橫條圖單色排序
BRAND_PALETTE = [
    BRAND_BLUE, BRAND_ORANGE, BRAND_YELLOW,
    BRAND_LBLUE, BRAND_BEIGE, BRAND_TEAL,
]


# ── 圖表圖例：緊貼圖面，並帶上件數 ─────────────────────────────
# 圖例與圖之間留太寬會把圖擠小；r 只留放得下文字的寬度。
PIE_LEGEND = dict(
    orientation="v",
    yanchor="middle", y=0.5,
    xanchor="left", x=1.0,
    font=dict(size=12),
    itemsizing="constant",
    tracegroupgap=2,
    bgcolor="rgba(0,0,0,0)",
    borderwidth=0,
)
PIE_MARGIN = dict(t=50, b=10, l=10, r=130)
BAR_LEGEND = dict(
    orientation="v",
    yanchor="top", y=1.0,
    xanchor="left", x=1.0,
    font=dict(size=12),
    tracegroupgap=2,
    bgcolor="rgba(0,0,0,0)",
    borderwidth=0,
)


def pie_legend_labels(counts) -> list[str]:
    """圓餅圖圖例：「名稱 n件」。

    圓餅圖整張只有一個 trace，圖例讀的是每個扇形的 label，
    所以件數必須寫進 names，改 trace.name 是沒有作用的。
    """
    return [f"{k} {int(v)}件" for k, v in counts.items()]


def add_counts_to_legend(fig, counts) -> None:
    """長條圖等「一個系列一個 trace」的圖：把件數接在系列名稱後面。"""
    mapping = {str(k): int(v) for k, v in counts.items()}

    def _rename(tr):
        n = mapping.get(str(tr.name))
        if n is not None:
            tr.update(name=f"{tr.name} {n}件")

    fig.for_each_trace(_rename)


# ── 啟動時確保 CJK 字型可用（下載備援）──────────────────────────
@st.cache_resource(show_spinner=False)
def _ensure_cjk_font() -> str:
    """回傳可用的 CJK 字型路徑；若系統沒有則下載到 /tmp。"""
    import os, glob
    CANDIDATES = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Medium.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/arphic/uming.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/tmp/NotoSansCJK.ttc",
    ]
    CANDIDATES += glob.glob("/usr/share/fonts/**/NotoSansCJK*.ttc", recursive=True)
    found = next((p for p in CANDIDATES if os.path.exists(p)), None)
    if found:
        return found
    # 下載到 /tmp
    _dl = "/tmp/NotoSansCJK.ttc"
    URLS = [
        "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTC/NotoSansCJK-Regular.ttc",
        "https://github.com/notofonts/noto-cjk/raw/main/Sans/OTC/NotoSansCJK-Regular.ttc",
    ]
    for url in URLS:
        try:
            import urllib.request
            urllib.request.urlretrieve(url, _dl)
            if os.path.exists(_dl) and os.path.getsize(_dl) > 100_000:
                return _dl
        except Exception:
            continue
    return ""

HISTORY_DIR = Path("history_reports")
HISTORY_DIR.mkdir(exist_ok=True)
META_FILE = HISTORY_DIR / "history.json"
DEFAULT_HISTORY_SHEET_ID = "1Sqh_8bXtFw7jvmCPufTpStKxfIafDzwYJRlgc0HFBSs"
SHEET_CELL_CHAR_LIMIT = 49000

# 範本路徑：優先使用與程式同目錄的 簡報範本.pptx（已隨程式一起部署）
TEMPLATE_PATH = Path(__file__).parent / "簡報範本.pptx"




def apply_brand_theme() -> None:
    st.markdown(
        """
        <style>
          html, body, [data-testid="stAppViewContainer"] {
            font-size: 18px !important;
          }
          @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@500;700;900&display=swap');
          
          /* Noto Sans TC Medium (500) - scoped to app content only, not Streamlit portals */
          [data-testid="stAppViewContainer"] *:not(.stIconMaterial):not(.material-symbols-rounded):not([data-testid="stIconMaterial"]),
          [data-testid="stHeader"] *:not(.stIconMaterial):not(.material-symbols-rounded):not([data-testid="stIconMaterial"]),
          [data-testid="stMain"] *:not(.stIconMaterial):not(.material-symbols-rounded):not([data-testid="stIconMaterial"]) {
            font-family: 'Noto Sans TC', 'Microsoft JhengHei', sans-serif !important;
          }
          [data-testid="stAppViewContainer"] p,
          [data-testid="stAppViewContainer"] span,
          [data-testid="stAppViewContainer"] label,
          [data-testid="stAppViewContainer"] div {
            font-weight: 500;
            font-size: 18px !important;
          }
          
          /* Use Noto Sans TC Medium (500) for everything — no bold allowed */
          h1, h2, h3, h4, h5, h6, .ecoco-banner, strong, b, .side-title, section[data-testid="stSidebar"] .stButton > button {
            font-family: 'Noto Sans TC', 'Microsoft JhengHei', sans-serif !important;
            font-weight: 500 !important;
          }

          /* ECOCO VI 品牌色票（Pantone 對應），介面只用這幾色 */
          :root{
            --ecoco-orange:#FF5000;      /* Pantone Orange 021 C */
            --ecoco-blue:#060E9F;        /* Pantone Blue 072 C   */
            --ecoco-yellow:#FFCE00;      /* Pantone 116 C        */
            --ecoco-lightblue:#8EB9C9;   /* Pantone 550 C        */
            --ecoco-beige:#FAE0B8;       /* Pantone P17-2 C      */
            --ecoco-deepteal:#0076A9;    /* Pantone 7690 C       */
            --ecoco-white:#FFFFFF;       /* Pantone White C      */
            /* 由品牌色衍生的介面用色：內文用深藍、次要文字用 7690、
               分隔線與淺底用 550 的淡化版，不引入品牌外的灰階。 */
            --ecoco-text:#060E9F;
            --ecoco-text-muted:#0076A9;
            --ecoco-line:rgba(142,185,201,.55);
            --ecoco-surface:rgba(142,185,201,.12);
          }
          .stApp {background: var(--ecoco-white);}

          /* ── 主標題：頂部橫條左上角 ─────────────────────────────
             標題放在 Streamlit 的 header 區，才能像截圖那樣橫跨整個寬度、
             並且待在側邊欄上方。 */
          [data-testid="stHeader"] {
            background: #FFFFFF !important;
            border-bottom: 2px solid var(--ecoco-orange);
            height: 62px !important;
          }
          [data-testid="stHeader"]::before {
            content: "ECOCO 客訴智能分析平台";
            position: absolute;
            left: 22px; top: 50%;
            transform: translateY(-50%);
            font-family: 'Noto Sans TC', 'Microsoft JhengHei', sans-serif;
            font-size: 30px;
            font-weight: 700;
            color: var(--ecoco-blue);
            letter-spacing: 0.5px;
            white-space: nowrap;
          }

          /* ── 頁首：淺底大標 + 小字副標，不加色塊 ────────────────── */
          .page-header {
            background: none;
            border: none;
            padding: 0;
            margin: 2px 0 22px;
          }
          [data-testid="stAppViewContainer"] .page-header .page-header-title,
          .page-header .page-header-title {
            font-size: 28px !important;
            font-weight: 700 !important;
            color: var(--ecoco-blue) !important;
            letter-spacing: 0.3px;
            margin: 0;
            line-height: 1.25;
          }
          [data-testid="stAppViewContainer"] .page-header .page-header-sub,
          .page-header .page-header-sub {
            font-size: 15px !important;
            font-weight: 500 !important;
            color: var(--ecoco-text-muted) !important;
            margin: 4px 0 0;
          }

          /* ── 首頁功能卡（2×2 網格）───────────────────────────── */
          .home-grid {
            display: grid;
            grid-template-columns: repeat(2, minmax(0, 1fr));
            gap: 18px;
          }
          @media (max-width: 900px) {
            .home-grid { grid-template-columns: minmax(0, 1fr); }
          }
          .home-card {
            background: var(--ecoco-white);
            border: 1px solid var(--ecoco-line);
            border-top: 4px solid var(--ecoco-orange);
            border-radius: 14px;
            padding: 22px 24px;
            display: flex;
            gap: 18px;
            align-items: flex-start;
          }
          .home-card-icon {
            flex: 0 0 auto;
            width: 46px; height: 46px;
            color: var(--ecoco-blue);
          }
          .home-card-icon svg { width: 46px; height: 46px; }
          [data-testid="stAppViewContainer"] .home-card-title,
          .home-card-title {
            font-size: 22px !important;
            font-weight: 700 !important;
            color: var(--ecoco-blue) !important;
            margin: 0 0 10px;
            line-height: 1.3;
          }
          .home-card ul {
            margin: 0; padding-left: 20px;
          }
          [data-testid="stAppViewContainer"] .home-card li,
          .home-card li {
            font-size: 15px !important;
            font-weight: 500 !important;
            color: var(--ecoco-text-muted) !important;
            line-height: 1.75;
          }
          .home-card li b { color: var(--ecoco-blue) !important; font-weight: 700 !important; }

          .feature-title {
            color: var(--ecoco-blue) !important;
            font-size: 18px !important;
            font-weight: 700 !important;
            margin: 0 0 10px 0;
          }
          .ecoco-card{
            border:1px solid var(--ecoco-line); border-left:6px solid var(--ecoco-orange);
            border-radius:12px; padding:10px 14px; background:var(--ecoco-white);
            margin-bottom:10px;
            color: var(--ecoco-text-muted) !important;
          }
          [data-testid="stAppViewContainer"] .ecoco-card,
          [data-testid="stAppViewContainer"] .ecoco-card * {
            font-size: 16px !important;
          }
          .ecoco-card b { color: var(--ecoco-blue) !important; }
          .small-muted { color: var(--ecoco-text-muted) !important; font-size: 0.9rem; }

          /* ── 側邊欄導覽（Pantone Blue 072 C）────────────────────── */
          section[data-testid="stSidebar"] {
            background: var(--ecoco-blue);
            border-right: 3px solid var(--ecoco-orange);
            /* 26px 的中文頁籤在預設 300px 寬會折行，加寬讓每個頁籤都排得下一行 */
            width: 370px !important;
            min-width: 370px !important;
          }
          section[data-testid="stSidebar"] [data-testid="stSidebarUserContent"] {
            padding-top: 6px;
          }
          /* 收合鈕不是導覽項目，不要套用下面的樣式，也不需要佔一整列 */
          section[data-testid="stSidebar"] [data-testid="stSidebarHeader"] {
            padding: 0 8px;
            min-height: 0;
          }
          section[data-testid="stSidebar"] [data-testid="stSidebarCollapseButton"] button {
            color: var(--ecoco-white) !important;
            background: transparent !important;
            border: none !important;
            min-height: 0 !important;
          }
          section[data-testid="stSidebar"] hr {
            border: none;
            border-top: 1px solid var(--ecoco-lightblue) !important;
            opacity: .6;
            margin: 12px 8px;
          }

          /* 導覽項目：整列可點、白字，選中/hover 才有底色 */
          section[data-testid="stSidebar"] .stButton > button {
            background-color: transparent !important;
            border: none !important;
            border-left: 4px solid transparent !important;
            border-radius: 0 8px 8px 0;
            color: var(--ecoco-white) !important;
            font-size: 26px !important;
            font-weight: 500 !important;
            line-height: 1.25;
            text-align: left;
            justify-content: flex-start;
            min-height: 56px;
            padding: 10px 14px;
            white-space: normal;
            transition: background-color .12s ease, color .12s ease !important;
          }
          section[data-testid="stSidebar"] .stButton > button * {
            color: var(--ecoco-white) !important;
            font-size: 26px !important;
            text-align: left;
          }
          section[data-testid="stSidebar"] .stButton > button:hover,
          section[data-testid="stSidebar"] .stButton > button:focus,
          section[data-testid="stSidebar"] .stButton > button:focus-visible,
          section[data-testid="stSidebar"] .stButton > button:active,
          section[data-testid="stSidebar"] .stButton > button[kind="primary"] {
            background-color: var(--ecoco-orange) !important;
            border-left: 4px solid var(--ecoco-yellow) !important;
            color: var(--ecoco-white) !important;
          }
          section[data-testid="stSidebar"] .stButton > button:hover *,
          section[data-testid="stSidebar"] .stButton > button:focus *,
          section[data-testid="stSidebar"] .stButton > button:focus-visible *,
          section[data-testid="stSidebar"] .stButton > button:active *,
          section[data-testid="stSidebar"] .stButton > button[kind="primary"] * {
            color: var(--ecoco-white) !important;
          }

          /* ── 表格工具列（欄位管理 / 自訂選項 / 批次問題處理）────── */
          .toolbar-title {
            font-size: 15px !important;
            font-weight: 700 !important;
            color: var(--ecoco-blue) !important;
            margin: 0 0 2px;
          }
          .st-key-editor_toolbar {
            background: var(--ecoco-surface);
            border: 1px solid var(--ecoco-line);
            border-radius: 12px;
            padding: 14px 16px 6px;
            margin-bottom: 12px;
          }
          .st-key-editor_toolbar [data-testid="stTextInput"] input,
          .st-key-editor_toolbar [data-testid="stSelectbox"] div[data-baseweb="select"] > div {
            font-size: 15px !important;
          }
          .st-key-editor_toolbar .stButton > button {
            min-height: 40px;
            border-radius: 8px;
            font-size: 15px !important;
            font-weight: 700 !important;
          }
          /* 主要動作用品牌橘，破壞性動作用深藍（品牌色內沒有紅色） */
          .st-key-editor_toolbar .stButton > button[kind="primary"] {
            background-color: var(--ecoco-orange) !important;
            border-color: var(--ecoco-orange) !important;
            color: var(--ecoco-white) !important;
          }
          .st-key-editor_toolbar .stButton > button[kind="secondary"] {
            background-color: var(--ecoco-white) !important;
            border: 1.5px solid var(--ecoco-blue) !important;
            color: var(--ecoco-blue) !important;
          }
          .st-key-editor_toolbar .stButton > button[kind="secondary"]:hover {
            background-color: var(--ecoco-blue) !important;
            color: var(--ecoco-white) !important;
          }

          /* Thicker scrollbar */
          ::-webkit-scrollbar { width: 10px; height: 10px; }
          ::-webkit-scrollbar-track { background: var(--ecoco-surface); border-radius: 6px; }
          ::-webkit-scrollbar-thumb { background: var(--ecoco-lightblue); border-radius: 6px; }
          ::-webkit-scrollbar-thumb:hover { background: var(--ecoco-blue); }

          /* File badge */
          .file-badge {
            display:inline-block; max-width:100%; padding:3px 10px;
            background:var(--ecoco-surface); border:1px solid var(--ecoco-lightblue);
            border-radius:20px;
            font-size:0.82rem; color:var(--ecoco-blue); white-space:nowrap;
            overflow:hidden; text-overflow:ellipsis; vertical-align:middle;
          }
          .editor-toolbar-title {
            font-size: 15px !important;
            font-weight: 700 !important;
            color: var(--ecoco-blue) !important;
            margin: 6px 0 4px;
          }
          [data-testid="stDataFrame"], [data-testid="stDataEditor"] {
            border: 1.5px solid var(--ecoco-lightblue) !important;
            border-radius: 6px !important;
            overflow-x: auto !important;
          }

          /* 移除 arrow_down 及內建圖示，避免異常顯示純文字 */
          [data-testid="stExpanderToggleIcon"], .material-symbols-rounded {
              display: none !important;
          }

        </style>
        """,
        unsafe_allow_html=True,
    )



def parse_pdf_to_df(file_obj) -> pd.DataFrame:
    if pdfplumber is None:
        raise RuntimeError("未安裝 pdfplumber，無法解析 PDF。")
    rows: list[dict] = []
    with pdfplumber.open(file_obj) as pdf:
        for p_idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            for ln_idx, line in enumerate(text.splitlines(), start=1):
                cleaned = re.sub(r"\s+", " ", line).strip()
                if cleaned:
                    rows.append({"page": p_idx, "line": ln_idx, "content": cleaned})
    return pd.DataFrame(rows if rows else [{"content": ""}])


def load_input_file(uploaded_file, filename: str = "") -> pd.DataFrame:
    """Load file from a Streamlit UploadedFile or BytesIO. Pass filename when using BytesIO."""
    name = filename or getattr(uploaded_file, "name", "")
    suffix = Path(name).suffix.lower()
    if suffix in [".xlsx", ".xls"]:
        return pd.read_excel(uploaded_file)
    if suffix == ".csv":
        for enc in ["utf-8-sig", "utf-8", "cp950", "big5"]:
            try:
                uploaded_file.seek(0)
                return pd.read_csv(uploaded_file, encoding=enc)
            except (UnicodeDecodeError, AttributeError):
                continue
        uploaded_file.seek(0)
        return pd.read_csv(uploaded_file, encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return parse_pdf_to_df(uploaded_file)
    raise ValueError(f"僅支援 excel / csv / pdf，收到：{suffix or name}")




@st.cache_resource(show_spinner="正在讀取歷史標記，建立分類知識庫…")
def get_knowledge(version: str = ""):
    """從歷史紀錄建立知識庫（指紋快取／挖掘規則／相似案例投票池）。

    version 只是 cache key：存檔後傳入新值即可讓知識庫重建。
    讀不到歷史時回傳 None，分類自動退回內建規則。
    """
    try:
        from automation.pipeline import build_knowledge

        return build_knowledge()   # 同時學習雲端歷史紀錄與本機 history_reports/
    except Exception as exc:       # 知識庫是加分項，失敗不能擋住分析
        st.session_state["_knowledge_error"] = str(exc)[:300]
        return None


def bump_knowledge_version() -> None:
    """人工修正存檔後呼叫，讓下次取用時重建知識庫（回饋閉環）。"""
    st.session_state["_knowledge_version"] = datetime.now().strftime("%Y%m%d%H%M%S")


def get_classifier():
    """組出分類器（依設定決定是否掛上選配的 LLM 層）。"""
    return build_classifier(
        knowledge=get_knowledge(st.session_state.get("_knowledge_version", ""))
    )


def analyze_dataframe(df: pd.DataFrame, cfg: AnalysisConfig,
                      classifier=None, progress=None) -> pd.DataFrame:
    """分析核心已移至 automation/core.py，此處保留同名包裝維持既有呼叫方式。"""
    if classifier is None:
        classifier = get_classifier()
    return core_analyze_dataframe(df, cfg, classifier=classifier, progress=progress)



# ── Google Sheets 歷史紀錄持久化 ────────────────────────────────────────────
# Render 的磁碟每次重啟會清空；使用 Google Sheets 作為永久儲存後端。
# 需在 Streamlit Secrets 設定：
#   HISTORY_SHEET_ID = "<your_spreadsheet_id>"
#   [google_credentials]   ← service account JSON 欄位

def _get_gsheet_client():
    """從環境變數或 st.secrets 取得 gspread client。"""
    try:
        import gspread as _gs
        from google.oauth2.service_account import Credentials as _Creds
    except ImportError:
        return None
    try:
        # 憑證來源統一由 automation/config.py 解析（環境變數 → secrets → JSON 檔）
        creds_dict = auto_config.get_google_credentials()
        if not creds_dict:
            return None
        creds = _Creds.from_service_account_info(
            creds_dict,
            scopes=["https://spreadsheets.google.com/feeds",
                    "https://www.googleapis.com/auth/drive"],
        )
        return _gs.authorize(creds)
    except Exception:
        return None


def _history_sheet(log_error: bool = False):
    """回傳歷史紀錄工作表；失敗回傳 None。log_error=True 時把錯誤存入 session_state。"""
    import os
    client = _get_gsheet_client()
    if client is None:
        if log_error:
            st.session_state["_gsheet_error"] = "無法建立 Google API 連線（請確認 GOOGLE_CREDENTIALS_JSON 環境變數格式正確）"
        return None
    try:
        sid = os.environ.get("HISTORY_SHEET_ID", "").strip()
        if not sid:
            try:
                sid = str(st.secrets.get("HISTORY_SHEET_ID", "")).strip()
            except Exception:
                sid = ""
        if not sid:
            sid = DEFAULT_HISTORY_SHEET_ID
        if not sid:
            if log_error:
                st.session_state["_gsheet_error"] = "未設定 HISTORY_SHEET_ID 環境變數"
            return None
        ss = client.open_by_key(sid)
        try:
            ws = ss.worksheet("歷史紀錄")
            try:
                header = ws.row_values(1)
                if header[:5] != ["id", "created_at", "source_name", "rows", "data_ref"]:
                    ws.update(values=[["id", "created_at", "source_name", "rows", "data_ref"]], range_name="A1:E1")
            except Exception:
                pass
            st.session_state.pop("_gsheet_error", None)
            return ws
        except Exception:
            ws = ss.add_worksheet("歷史紀錄", rows=500, cols=6)
            ws.append_row(["id", "created_at", "source_name", "rows", "data_ref"])
            st.session_state.pop("_gsheet_error", None)
            return ws
    except Exception as e:
        err_str = str(e)
        if "PERMISSION_DENIED" in err_str or "403" in err_str:
            msg = (f"Google Sheets API 權限錯誤。請至 Google Cloud Console 確認已啟用：\n"
                   f"1. Google Sheets API\n2. Google Drive API\n"
                   f"錯誤：{err_str[:200]}")
        elif "NOT_FOUND" in err_str or "404" in err_str:
            msg = f"試算表不存在（ID 可能錯誤）：{err_str[:200]}"
        else:
            msg = f"Google Sheets 連線錯誤：{err_str[:300]}"
        if log_error:
            st.session_state["_gsheet_error"] = msg
        return None


def _sanitize_sheet_value(value, max_chars: int = SHEET_CELL_CHAR_LIMIT) -> str:
    if pd.isna(value):
        return ""
    text = str(value)
    if text.lower() in {"nan", "inf", "-inf", "infinity", "-infinity"}:
        return ""
    return text[:max_chars]


def _sanitize_df_for_sheet(df: pd.DataFrame, max_chars: int = SHEET_CELL_CHAR_LIMIT) -> pd.DataFrame:
    out = df.copy()
    out = out.replace([float("inf"), float("-inf")], pd.NA)
    out = out.astype(object).where(pd.notna(out), "")
    mapper = lambda v: _sanitize_sheet_value(v, max_chars=max_chars)
    if hasattr(out, "map"):
        return out.map(mapper)
    return out.apply(lambda col: col.map(mapper))


def _history_data_sheet_name(item_id: str) -> str:
    safe_id = re.sub(r"[^0-9A-Za-z_\\-]+", "_", str(item_id))[:80]
    return f"history_{safe_id}"


def _write_history_data_sheet(spreadsheet, worksheet_name: str, df: pd.DataFrame):
    clean_df = _sanitize_df_for_sheet(df)
    values = [clean_df.columns.tolist()] + clean_df.values.tolist()
    rows = max(len(values), 1)
    cols = max(len(clean_df.columns), 1)
    try:
        ws_data = spreadsheet.worksheet(worksheet_name)
        ws_data.clear()
        ws_data.resize(rows=max(rows, 100), cols=max(cols, 10))
    except Exception:
        ws_data = spreadsheet.add_worksheet(title=worksheet_name, rows=max(rows, 100), cols=max(cols, 10))
    if values:
        ws_data.update(values=values, range_name="A1")
    return ws_data


def _worksheet_to_dataframe(ws) -> pd.DataFrame:
    values = ws.get_all_values()
    if not values:
        return pd.DataFrame()
    header = values[0]
    rows = values[1:]
    width = len(header)
    normalized = [(row + [""] * width)[:width] for row in rows]
    return pd.DataFrame(normalized, columns=header)


def save_history(df: pd.DataFrame, source_name: str, existing_id: str = "") -> tuple[Path, str, str]:
    if auto_config.history_readonly():
        # 唯讀模式（測試／示範用）：不寫雲端也不寫本機，避免污染正式歷史
        st.session_state["_gsheet_error"] = "唯讀模式（HISTORY_READONLY=true），未寫入歷史紀錄"
        return Path(), "", existing_id
    today = datetime.now().strftime("%Y%m%d")
    ts = existing_id if existing_id else datetime.now().strftime("%Y%m%d_%H%M%S")
    output_name = f"{today}_分析.xlsx"
    excel_bytes = to_excel_bytes(df)
    data_sheet_name = _history_data_sheet_name(ts)

    meta = {
        "id": ts, "created_at": datetime.now().isoformat(timespec="seconds"),
        "source_name": source_name, "output_name": output_name,
        "output_path": "", "rows": int(len(df)),
    }

    saved_to_gsheet = False
    # 1. Google Sheets（永久）
    ws = _history_sheet(log_error=True)
    if ws is not None:
        try:
            _write_history_data_sheet(ws.spreadsheet, data_sheet_name, df)
            if existing_id:
                rows = ws.get_all_values()
                for i, row in enumerate(rows[1:], start=2):
                    if row and row[0] == existing_id:
                        ws.delete_rows(i); break
            ws.append_row([ts, meta["created_at"], source_name, str(len(df)), f"sheet:{data_sheet_name}"])
            st.session_state.pop("_gsheet_error", None)
            saved_to_gsheet = True
        except Exception as e:
            st.session_state["_gsheet_error"] = f"歷史紀錄寫入 Google Sheets 失敗：{str(e)[:300]}"

    if saved_to_gsheet:
        if "_history_cache" not in st.session_state:
            st.session_state["_history_cache"] = {}
        st.session_state["_history_cache"][ts] = {"meta": meta, "excel_bytes": excel_bytes}

    # 2. 本機磁碟（只保存檔案，不作為歷史紀錄來源）
    output_path = HISTORY_DIR / f"{ts}_{output_name}"
    try:
        output_path.write_bytes(excel_bytes)
    except Exception:
        pass
    return output_path, output_name, ts


def load_history() -> list[dict]:
    import base64
    merged: dict[str, dict] = {}

    # 只從 Google Sheets 讀取歷史紀錄，避免雲端無紀錄但網頁殘留。
    ws = _history_sheet()
    if ws is None:
        st.session_state["_history_cache"] = {}
        return []
    try:
        for row in ws.get_all_values()[1:]:
            if not row or not row[0]:
                continue
            rid = row[0]
            created_at = row[1] if len(row) > 1 else ""
            sname = row[2] if len(row) > 2 else ""
            rows_str = row[3] if len(row) > 3 else "0"
            data_ref = row[4] if len(row) > 4 else ""
            meta = {
                "id": rid, "created_at": created_at,
                "source_name": sname,
                "rows": int(rows_str) if rows_str.isdigit() else 0,
                "output_name": f"{rid}_分析.xlsx", "output_path": "",
            }
            merged[rid] = meta
            if "_history_cache" not in st.session_state:
                st.session_state["_history_cache"] = {}
            if rid not in st.session_state["_history_cache"] and data_ref:
                try:
                    if data_ref.startswith("sheet:"):
                        data_ws = ws.spreadsheet.worksheet(data_ref.split(":", 1)[1])
                        # 遮蔽功能上線前存的舊紀錄，讀回來時補遮一次
                        hist_df = mask_phone_columns(_worksheet_to_dataframe(data_ws))
                        excel_bytes = to_excel_bytes(hist_df)
                    else:
                        excel_bytes = base64.b64decode(data_ref)
                    st.session_state["_history_cache"][rid] = {
                        "meta": meta,
                        "excel_bytes": excel_bytes,
                    }
                except Exception:
                    pass
    except Exception:
        st.session_state["_history_cache"] = {}
        return []

    return sorted(merged.values(), key=lambda x: x.get("created_at", ""), reverse=True)


def safe_filename(text: str) -> str:
    return re.sub(r'[\\/:*?"<>|]+', "_", str(text))


def delete_history(item_id: str):
    ws = _history_sheet()
    if ws:
        try:
            for i, row in enumerate(ws.get_all_values()[1:], start=2):
                if row and row[0] == item_id:
                    ws.delete_rows(i); break
        except Exception:
            pass
    if META_FILE.exists():
        try:
            history = json.loads(META_FILE.read_text(encoding="utf-8"))
            history = [i for i in history if i["id"] != item_id]
            META_FILE.write_text(json.dumps(history, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception:
            pass
    cache = st.session_state.get("_history_cache", {})
    cache.pop(item_id, None)
    st.session_state["_history_cache"] = cache



def generate_ai_summary(df: pd.DataFrame) -> str:
    if df.empty:
        return "目前沒有可分析資料。"
    total = len(df)
    type_count = df["問題類型"].value_counts()
    detail_count = df["問題細項"].value_counts()
    top_type = type_count.index[0]
    top_type_count = int(type_count.iloc[0])
    top_detail = detail_count.index[0]
    top_detail_count = int(detail_count.iloc[0])
    dept_text = ""
    if "部門" in df.columns and not df["部門"].dropna().empty:
        dept_top = df["部門"].replace("", "未分配").value_counts().head(3)
        dept_text = "；".join([f"{k} {int(v)} 件" for k, v in dept_top.items()])
    detail_lines = []
    for name, count in detail_count.head(5).items():
        detail_lines.append(f"{name} {int(count)} 件")
    detail_text = "；".join(detail_lines)
    return (
        f"1) 自動摘要：本次共 {total} 件，主力問題為「{top_type}」{top_type_count} 件，占比 {top_type_count/total:.1%}。\n"
        f"2) 細項說明：最高頻細項為「{top_detail}」{top_detail_count} 件；TOP5 為 {detail_text}。\n"
        f"3) 部門觀察：{dept_text or '目前無明確部門欄位可判讀'}。\n"
        "4) 初步判讀：請優先檢查高頻細項是否集中於特定站點、設備型態或操作流程，並比對近期是否有維修、活動或系統異動。\n"
        "5) 建議行動：以 TOP3 問題建立改善任務，指定負責部門、預計完成日與每週追蹤指標。"
    )


def _ai_summary_cache_key(df: pd.DataFrame, tag: str = "") -> str:
    """以資料指紋當快取鍵：資料沒變就不重新呼叫 API。"""
    try:
        fingerprint = pd.util.hash_pandas_object(
            df[[c for c in ("問題類型", "問題細項", "部門") if c in df.columns]],
            index=False,
        ).sum()
    except Exception:
        fingerprint = len(df)
    return f"_ai_summary_{tag}_{len(df)}_{fingerprint}"


def get_ai_summary_cached(df: pd.DataFrame, tag: str = "") -> str:
    """取得 AI 摘要（同一份資料只會呼叫一次 API）。

    Streamlit 每次互動都會整頁重跑，若直接呼叫會每動一下就計費一次。
    """
    key = _ai_summary_cache_key(df, tag)
    if key not in st.session_state:
        with st.spinner("正在產生 AI 摘要…"):
            st.session_state[key] = generate_ai_summary_llm(df)
    return st.session_state[key]


def generate_ai_summary_llm(df: pd.DataFrame, model_name: str = "") -> str:
    """用 LLM 產生摘要；沒有 API key 或呼叫失敗時退回內建統計摘要。

    注意：這個函式會產生 API 費用，呼叫端請走 get_ai_summary_cached()。
    """
    if not (auto_config.get_anthropic_key() or auto_config.get_openai_key()):
        return generate_ai_summary(df)
    sample = df[["問題類型", "問題細項", "部門"]].head(300).to_dict(orient="records")
    payload = {
        "total_rows": len(df),
        "top_types": df["問題類型"].value_counts().head(6).to_dict(),
        "top_details": df["問題細項"].value_counts().head(10).to_dict(),
        "sample_rows": sample,
    }
    prompt = (
        "你是客服品質分析顧問。請用繁體中文輸出3-5點重點，格式精簡，"
        "包含: 高頻問題、可能根因、跨部門優先改善建議。資料如下:\n"
        f"{json.dumps(payload, ensure_ascii=False)}\n"
        "請特別列出：1. 站點城市分布熱點 (如果從內容看得出來) 2. 問題類型與細項的熱點(最高頻的異常)。"
    )
    from automation.llm import complete_text

    text = complete_text(prompt, model=model_name or None, max_tokens=1200)
    return text.strip() if text else generate_ai_summary(df)


def _drop_ui_columns(df: pd.DataFrame) -> pd.DataFrame:
    """「選取」只是表格上的勾選狀態，不該出現在下載檔裡。"""
    return df.drop(columns=[c for c in ["選取"] if c in df.columns], errors="ignore")


def to_excel_bytes(df: pd.DataFrame) -> bytes:
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        _drop_ui_columns(df).to_excel(writer, index=False, sheet_name="analysis")
    return buffer.getvalue()


def to_csv_bytes(df: pd.DataFrame) -> bytes:
    return _drop_ui_columns(df).to_csv(index=False).encode("utf-8-sig")


def to_pdf_bytes(df: pd.DataFrame, source_name: str = "", download_count: int = 1) -> bytes:
    """Generate PDF using fpdf2 + Noto CJK for Traditional Chinese support."""
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    import os, glob

    # ── 找字型：多重備援路徑 ──
    CJK_FONT_CANDIDATES = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Medium.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJKtc-Regular.otf",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJKtc-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/arphic/uming.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    ]
    CJK_FONT_CANDIDATES += glob.glob("/usr/share/fonts/**/NotoSansCJK*.ttc", recursive=True)
    CJK_FONT_CANDIDATES += glob.glob("/usr/share/fonts/**/NotoSansCJK*.otf", recursive=True)
    font_path = next((p for p in CJK_FONT_CANDIDATES if os.path.exists(p)), None)
    # 使用 _ensure_cjk_font 確保有可用字型
    if not font_path:
        font_path = _ensure_cjk_font()

    table_df = df.copy()
    drop_cols = [c for c in ["選取"] if c in table_df.columns]
    table_df = table_df.drop(columns=drop_cols).fillna("")

    PAGE_W_MM = 277.0
    WIDE_COLS   = {"用戶內容", "主旨", "問題主旨"}
    MEDIUM_COLS = {"問題細項", "問題類型", "進件日期", "日期時間", "站點名稱", "問題細項"}
    num_cols = len(table_df.columns)
    wide_count   = sum(1 for c in table_df.columns if c in WIDE_COLS)
    medium_count = sum(1 for c in table_df.columns if c in MEDIUM_COLS)
    narrow_count = num_cols - wide_count - medium_count
    unit = PAGE_W_MM / max(wide_count*4 + medium_count*2 + narrow_count, 1)
    col_widths = {}
    for c in table_df.columns:
        if c in WIDE_COLS:       col_widths[c] = unit * 4
        elif c in MEDIUM_COLS:   col_widths[c] = unit * 2
        else:                    col_widths[c] = unit

    pdf = FPDF(orientation="L", format="A4")
    pdf.set_auto_page_break(auto=True, margin=10)
    pdf.add_page()

    if font_path:
        try:
            pdf.add_font("CJK", style="", fname=font_path)
            FONT = "CJK"
        except Exception:
            FONT = "Helvetica"
    else:
        FONT = "Helvetica"

    ROW_H = 7.0; HDR_H = 8.0; FS_HDR = 8; FS_CELL = 7

    def safe_text(s):
        if FONT == "Helvetica":
            return s.encode("ascii", "replace").decode()
        return s

    def fit_text_in_cell(pdf_obj, text, max_width, max_size=8, min_size=5):
        """縮小字型直到文字適合欄寬"""
        for fs in range(max_size, min_size-1, -1):
            pdf_obj.set_font(FONT, size=fs)
            if pdf_obj.get_string_width(text) <= max_width - 1:
                return fs
        return min_size

    def draw_page_label():
        label = f"{source_name or '分析檔案'}  {datetime.now().strftime('%Y/%m/%d')}  第 {download_count} 次"
        pdf.set_xy(8, 4)
        pdf.set_text_color(80, 80, 80)
        pdf.set_font(FONT, size=7)
        pdf.cell(0, 5, safe_text(label), align="L")
        pdf.set_xy(10, 10)

    draw_page_label()

    # 表頭（自動縮小以適應欄寬）
    pdf.set_fill_color(0x06, 0x0E, 0x9F)
    pdf.set_text_color(255, 255, 255)
    for col in table_df.columns:
        cw = col_widths[col]
        header_text = safe_text(col)
        fs = fit_text_in_cell(pdf, header_text, cw, max_size=FS_HDR, min_size=4)
        pdf.set_font(FONT, size=fs)
        pdf.cell(cw, HDR_H, header_text, border=1, fill=True,
                 new_x=XPos.RIGHT, new_y=YPos.TOP, align="C")
    pdf.ln(HDR_H)

    # ── 資料列：框線統一 row_h 高度，文字疊加換行 ──
    pdf.set_font(FONT, size=FS_CELL)
    col_list = list(table_df.columns)

    for i, (_, row) in enumerate(table_df.iterrows()):
        fill_rgb = (0xEB, 0xF4, 0xFA) if i % 2 == 0 else (0xFF, 0xFF, 0xFF)

        # ── 準備文字（去除多餘空白與換行）──
        cell_texts = {
            col: safe_text(
                " ".join(str(row[col]).split())   # 多空白合併為單一空格
                .replace("  ", " ").strip()
            )
            for col in col_list
        }

        # ── 精確計算每欄行數 ──
        col_lines: dict[str, int] = {}
        for col in col_list:
            cw = col_widths[col] - 2
            text = cell_texts[col]
            if not text:
                col_lines[col] = 1
                continue
            n = 0
            for para in text.replace("\r", "").split("\n"):
                if not para:
                    n += 1
                    continue
                line_w = 0.0
                for ch in para:
                    ch_w = pdf.get_string_width(ch)
                    if line_w + ch_w > cw:
                        n += 1
                        line_w = ch_w
                    else:
                        line_w += ch_w
                n += 1
            col_lines[col] = max(1, n)

        max_lines = max(col_lines.values())
        row_h = max_lines * ROW_H

        # ── 換頁檢查 ──
        if pdf.get_y() + row_h > pdf.page_break_trigger:
            pdf.add_page()
            draw_page_label()
            pdf.set_fill_color(0x06, 0x0E, 0x9F)
            pdf.set_text_color(255, 255, 255)
            for col in col_list:
                cw = col_widths[col]
                hdr_t = safe_text(col)
                fs = fit_text_in_cell(pdf, hdr_t, cw, max_size=FS_HDR, min_size=4)
                pdf.set_font(FONT, size=fs)
                pdf.cell(cw, HDR_H, hdr_t, border=1, fill=True,
                         new_x=XPos.RIGHT, new_y=YPos.TOP, align="C")
            pdf.ln(HDR_H)
            pdf.set_font(FONT, size=FS_CELL)

        x0 = pdf.get_x()
        y0 = pdf.get_y()
        pdf.set_text_color(0x22, 0x22, 0x22)

        # ── Step 1：先畫每欄的底色 + 完整框線（統一 row_h）──
        x_cursor = x0
        for col in col_list:
            cw = col_widths[col]
            # 底色填滿整格
            pdf.set_fill_color(*fill_rgb)
            pdf.rect(x_cursor, y0, cw, row_h, style="F")
            # 外框線（整格高度）
            pdf.set_draw_color(0x99, 0x99, 0x99)
            pdf.rect(x_cursor, y0, cw, row_h, style="D")
            x_cursor += cw

        # ── Step 2：疊加文字（multi_cell 只畫文字，不畫框線）──
        x_cursor = x0
        for col in col_list:
            cw = col_widths[col]
            val = cell_texts[col]
            pdf.set_xy(x_cursor + 0.5, y0 + 0.5)   # 0.5mm 內縮 padding
            pdf.set_fill_color(*fill_rgb)
            pdf.multi_cell(
                cw - 1, ROW_H, val,
                border=0,          # 不畫框線（已在 Step 1 畫好）
                align="L", fill=False,
                new_x=XPos.RIGHT, new_y=YPos.TOP,
                max_line_height=ROW_H,
            )
            x_cursor += cw

        pdf.set_draw_color(0, 0, 0)  # 恢復黑色
        pdf.set_xy(x0, y0 + row_h)

    # 頁尾
    pdf.set_y(-12)
    pdf.set_font(FONT, size=6)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 6, safe_text(f"ECOCO 客訴分析報告  共 {len(table_df)} 筆  產出日期：{datetime.now().strftime('%Y/%m/%d')}"), align="C")

    return bytes(pdf.output())

def _setup_cjk_font() -> None:
    """設定 matplotlib 中文字型，使用 _ensure_cjk_font 取得字型路徑。"""
    import matplotlib.font_manager as fm
    import os

    # ── 已設定過就直接返回 ──
    current = plt.rcParams.get("font.family", "")
    if current and "sans-serif" not in str(current) and current != ["DejaVu Sans"]:
        return

    # ── 優先使用 _ensure_cjk_font 確保字型存在 ──
    fp = _ensure_cjk_font()
    if fp and os.path.exists(fp):
        try:
            fm.fontManager.addfont(fp)
            plt.rcParams["font.family"] = fm.FontProperties(fname=fp).get_name()
            plt.rcParams["axes.unicode_minus"] = False
            return
        except Exception:
            pass

    # ── 1. 已知路徑（Ubuntu / Render / Debian）──
    KNOWN_PATHS = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Medium.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJKtc-Regular.otf",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJKtc-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/arphic/uming.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    ]
    for fp in KNOWN_PATHS:
        if os.path.exists(fp):
            try:
                fm.fontManager.addfont(fp)
                fname = fm.FontProperties(fname=fp).get_name()
                plt.rcParams["font.family"] = fname
                plt.rcParams["axes.unicode_minus"] = False
                return
            except Exception:
                continue

    # ── 2. 掃描所有已安裝字型找 CJK ──
    CJK_KEYWORDS = [
        "Noto Sans CJK", "Noto Serif CJK", "Noto Sans TC",
        "MingLiU", "PMingLiU", "Microsoft JhengHei",
        "WenQuanYi", "Droid Sans Fallback", "AR PL UMing",
    ]
    fm._load_fontmanager(try_read_cache=False)   # 強制重新掃描
    for kw in CJK_KEYWORDS:
        for f in fm.fontManager.ttflist:
            if kw.lower() in f.name.lower():
                plt.rcParams["font.family"] = f.name
                plt.rcParams["axes.unicode_minus"] = False
                return

    # ── 3. 最終 fallback：至少關掉負號亂碼 ──
    plt.rcParams["axes.unicode_minus"] = False


def build_chart_pack(df: pd.DataFrame,
                     color_bar: str | None = None,
                     color_pie: list[str] | None = None,
                     color_hbar: str | None = None) -> dict[str, bytes]:
    """Build chart PNG images for download/PPT.
    color_bar  : 問題類型直條圖 — None = 依部門品牌色; 或傳入單一 hex 強制套用
    color_pie  : 機台圓餅圖各扇形顏色 list，None = BRAND_PALETTE
    color_hbar : 十大細項橫條圖顏色，None = BRAND_BLUE
    """
    _setup_cjk_font()

    data = df.copy()
    # ── 機台類型正規化：方舟/方舟站 → 收瓶機 ──
    if "機台類型" in data.columns:
        data["機台類型"] = data["機台類型"].apply(
            lambda v: "收瓶機" if ("方舟" in str(v) or "收瓶" in str(v))
            else ("電池機" if "電池" in str(v) else str(v))
        )
    stats = data["問題類型"].value_counts().rename_axis("問題類型").reset_index(name="件數")
    stats["百分比"] = (stats["件數"] / max(stats["件數"].sum(), 1) * 100).round(1)
    detail_stats = data["問題細項"].value_counts().reset_index().head(10)
    detail_stats.columns = ["問題細項", "件數"]
    d = detail_stats.sort_values("件數", ascending=True)

    # ── resolve colors ──
    _pie_palette  = color_pie  if color_pie  else BRAND_PALETTE
    _hbar_color   = color_hbar if color_hbar else BRAND_BLUE

    def _bar_colors_for(series):
        if color_bar:
            return [color_bar] * len(series)
        return [DEPT_COLOR_MAP.get(DEPT_MAP.get(t, ""), BRAND_ORANGE) for t in series]

    # 1) 問題類型直條圖
    fig1, ax1 = plt.subplots(figsize=(8, 4.5))
    bc = _bar_colors_for(stats["問題類型"])
    ax1.bar(stats["問題類型"], stats["件數"], color=bc)
    ax1.set_title("問題類型分布")
    ax1.set_ylabel("件數")
    ax1.yaxis.set_major_locator(plt.MaxNLocator(integer=True))
    ax1.tick_params(axis="x", rotation=20)
    for i, r in stats.iterrows():
        ax1.text(i, r["件數"], f'{int(r["百分比"])}%', ha="center", va="bottom", fontsize=9)
    fig1.tight_layout()
    b1 = io.BytesIO(); fig1.savefig(b1, format="png", dpi=180); plt.close(fig1)

    # 2) 機台圓餅圖
    fig2, ax2 = plt.subplots(figsize=(6.2, 4.5))
    df_machine = data[data["問題類型"] == "機台問題類型"].copy()
    if df_machine.empty:
        ax2.text(0.5, 0.5, "無機台相關資料", ha="center", va="center", transform=ax2.transAxes)
        pie_counts = None
    else:
        def _get_mtype(row):
            txt = str(row.get("用戶內容", "")) + " " + str(row.get("主旨", ""))
            if "方舟" in txt: return "方舟站"
            if "電池" in txt: return "電池機"
            return "收瓶機"
        df_machine["機台機型"] = df_machine.apply(_get_mtype, axis=1)
        pie_counts = df_machine["機台機型"].value_counts()
        pc = _pie_palette[:len(pie_counts)]
        wedges, texts, autotexts = ax2.pie(
            pie_counts.values, labels=pie_counts.index, autopct="%1.1f%%",
            colors=pc, wedgeprops=dict(linewidth=1.5, edgecolor="white"),
        )
        for at in autotexts: at.set_fontsize(10)
    ax2.set_title("機台問題類型分布")
    fig2.tight_layout()
    b2 = io.BytesIO(); fig2.savefig(b2, format="png", dpi=180); plt.close(fig2)

    # 3) 十大細項橫條圖  ── 強制品牌主藍 #060E9F，整數刻度
    fig3, ax3 = plt.subplots(figsize=(8, 4.5))
    _hbar = _hbar_color if _hbar_color else "#060E9F"
    ax3.barh(d["問題細項"], d["件數"], color=_hbar)
    ax3.set_title("十大問題細項分布")
    ax3.set_xlabel("件數")
    # 強制整數刻度（件數必為整數）
    from matplotlib.ticker import MultipleLocator
    ax3.xaxis.set_major_locator(MultipleLocator(1))
    ax3.xaxis.set_minor_locator(MultipleLocator(1))
    ax3.set_xlim(left=0)
    fig3.tight_layout()
    b3 = io.BytesIO(); fig3.savefig(b3, format="png", dpi=180); plt.close(fig3)

    # 4) Dashboard 合圖
    fig4 = plt.figure(figsize=(14, 5))
    gs = fig4.add_gridspec(1, 3)
    a1 = fig4.add_subplot(gs[0, 0])
    a2 = fig4.add_subplot(gs[0, 1])
    a3 = fig4.add_subplot(gs[0, 2])
    a1.bar(stats["問題類型"], stats["件數"], color=bc)
    a1.set_title("問題類型分布")
    a1.yaxis.set_major_locator(MultipleLocator(1))
    a1.tick_params(axis="x", rotation=18)
    if pie_counts is None:
        a2.text(0.5, 0.5, "無機台資料", ha="center", va="center", transform=a2.transAxes)
    else:
        a2.pie(pie_counts.values, labels=pie_counts.index, autopct="%1.1f%%",
               colors=_pie_palette[:len(pie_counts)],
               wedgeprops=dict(linewidth=1.5, edgecolor="white"))
    a2.set_title("機台問題占比")
    a3.barh(d["問題細項"], d["件數"], color=_hbar)
    a3.xaxis.set_major_locator(MultipleLocator(1))
    a3.set_xlim(left=0)
    a3.set_title("十大細項")
    fig4.tight_layout()
    b4 = io.BytesIO(); fig4.savefig(b4, format="png", dpi=180); plt.close(fig4)

    return {
        "chart_問題類型分布.png": b1.getvalue(),
        "chart_機台問題占比.png": b2.getvalue(),
        "chart_十大問題細項.png": b3.getvalue(),
        "chart_dashboard.png":    b4.getvalue(),
    }


def build_ppt_bytes(stats: pd.DataFrame, ai_text: str, source_name: str,
                    template_path: str = "",
                    chart_pack: Optional[dict[str, bytes]] = None) -> bytes:
    """
    Build a PPT presentation.
    優先使用同目錄的範本；若找不到則從零構建符合 ECOCO 品牌風格的投影片。
    """
    from pptx.util import Emu, Inches, Pt
    from pptx.enum.text import PP_ALIGN

    BLUE   = RGBColor(0x06, 0x0E, 0x9F)
    ORANGE = RGBColor(0xFF, 0x50, 0x00)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    BEIGE  = RGBColor(0xFA, 0xE0, 0xB8)
    DARK   = RGBColor(0x22, 0x22, 0x22)
    LGRAY  = RGBColor(0xE8, 0xF1, 0xF5)
    FONT   = "MingLiU"   # 細明體

    # ── 嘗試載入範本 ──
    _tpath = Path(template_path) if template_path else TEMPLATE_PATH
    use_template = _tpath.exists()
    prs = Presentation(str(_tpath)) if use_template else Presentation()
    if not use_template:
        # 設定投影片大小為寬螢幕 16:9
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    SW = prs.slide_width
    SH = prs.slide_height

    # ── 小工具 ──
    def blank_layout():
        for lay in prs.slide_layouts:
            if lay.name.lower() in ("blank", "空白"):
                return lay
        return prs.slide_layouts[-1]

    def add_rect(slide, l, t, w, h, fill_rgb, line=False):
        shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
        if not line:
            shp.line.fill.background()
        return shp

    def add_text(slide, text, l, t, w, h, font_size, bold=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color or DARK
        return txb

    def add_img(slide, img_bytes, l, t, w, h):
        if img_bytes:
            try:
                slide.shapes.add_picture(io.BytesIO(img_bytes),
                    Inches(l), Inches(t), Inches(w), Inches(h))
            except Exception:
                pass

    def add_header(slide, title_text, subtitle_text=""):
        """加上 ECOCO 品牌頁首（藍色長條 + 標題）"""
        add_rect(slide, 0, 0, SW/914400, 1.05, BLUE)
        add_text(slide, title_text, 0.3, 0.08, 9.0, 0.55,
                 20, bold=True, color=WHITE)
        if subtitle_text:
            add_text(slide, subtitle_text, 0.3, 0.62, 10.0, 0.38,
                     11, color=BEIGE)

    def delete_shape(sp):
        sp.element.getparent().remove(sp.element)

    # ════════════════════════════════════════════════════════
    #  使用範本：覆寫文字、表格、圖片
    # ════════════════════════════════════════════════════════
    if use_template:
        slides = list(prs.slides)

        # --- 封面 (slide 0) ---
        # 範本版面：左側藍色面板（版面配置提供），右側三個文字框：
        #   Shape;99  → 主標題「營運周報」 (l≈5.66, t≈2.18)
        #   Shape;98  → 日期/資料列 (l≈6.14, t≈3.48)
        #   Shape;96  → 公司名藍底白字 (l≈6.67, t≈5.04)
        s0 = slides[0]
        for sp in s0.shapes:
            if not sp.has_text_frame:
                continue
            l_in = sp.left / 914400
            t_in = sp.top  / 914400
            raw  = sp.text_frame.text.strip()

            # ── 主標題（在 x>5" 且 y<3"）
            if l_in > 5.0 and t_in < 3.0:
                tf = sp.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "客訴分析簡報"
                run.font.name  = FONT
                run.font.bold  = True
                run.font.size  = Pt(32)
                run.font.color.rgb = RGBColor(0x16, 0x2B, 0x7E)

            # ── 日期/資料欄（在 x>5" 且 y 在 3~5"）
            elif l_in > 5.0 and 3.0 <= t_in < 5.0:
                tf = sp.text_frame
                tf.clear()
                for label, val in [
                    ("報告日期", datetime.now().strftime("%Y/%m/%d")),
                    ("報告資料", source_name),
                ]:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = f"{label}:{val}"
                    run.font.name  = FONT
                    run.font.bold  = True
                    run.font.size  = Pt(18)
                    run.font.color.rgb = RGBColor(0x1A, 0x2A, 0x7F)

            # ── 公司名（在 x>6" 且 y>=5" 或有填色藍底）
            elif l_in > 6.0 and t_in >= 4.8:
                pass   # 保留原樣「凡立橙股份有限公司」

        def _fill_slide(slide, title_txt, chart_key_list, add_table=True):
            SWi = prs.slide_width  / 914400
            SHi = prs.slide_height / 914400

            # 更新標題文字（比對關鍵字）
            for sp in slide.shapes:
                if sp.has_text_frame:
                    txt = sp.text_frame.text
                    if any(k in txt for k in ("客訴問題分析", "機台問題佔比", "機台與細項",
                                               "客訴問題", "問題分析", "20260")):
                        tf = sp.text_frame; tf.clear()
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = title_txt
                        run.font.name = FONT
                        run.font.bold = True
                        run.font.size = Pt(16)
                        run.font.color.rgb = BLUE

            # 收集現有 Table / Picture 位置後刪除（清空舊內容）
            tbl_rect = None
            pic_rects = []
            for sp in list(slide.shapes):
                if sp.shape_type == 19:   # Table
                    tbl_rect = (sp.left, sp.top, sp.width, sp.height)
                    delete_shape(sp)
                elif sp.shape_type == 13:  # Picture
                    pic_rects.append((sp.left, sp.top, sp.width, sp.height))
                    delete_shape(sp)
            pic_rects.sort(key=lambda x: x[0])

            # ── 圖表插入：優先使用範本佔位位置，否則用固定座標 ──
            if chart_pack:
                if add_table:
                    # slide 2（問題分析）：表格左半 + 圖表右半
                    # 固定座標：圖表放右側
                    chart_fixed = [
                        (6.2, 1.15, SWi - 6.5, SHi - 1.4),   # 問題類型分布
                    ]
                else:
                    # slide 3（機台細項）：左右各放一張圖
                    chart_fixed = [
                        (0.3,              1.15, (SWi - 0.6) / 2,       SHi - 1.4),
                        (0.3 + (SWi-0.6)/2 + 0.15, 1.15, (SWi-0.6)/2, SHi - 1.4),
                    ]

                for idx, key in enumerate(chart_key_list):
                    if key not in chart_pack:
                        continue
                    if idx < len(pic_rects):
                        # 範本有佔位圖片 → 用原始位置
                        add_img(slide, chart_pack[key],
                                *[v / 914400 for v in pic_rects[idx]])
                    elif idx < len(chart_fixed):
                        # 範本沒有佔位 → 用固定座標
                        add_img(slide, chart_pack[key], *chart_fixed[idx])

            # ── 重建資料表格 ──
            if add_table:
                # 如果範本有舊表格位置就沿用，否則預設左側
                if tbl_rect:
                    tb_l, tb_t, tb_w, tb_h = tbl_rect
                else:
                    tb_l = Inches(0.25)
                    tb_t = Inches(1.15)
                    tb_w = Inches(5.8)
                    tb_h = Inches(SHi - 1.4)
                rows_n = min(len(stats) + 1, 12)
                tb = slide.shapes.add_table(rows_n, 4, tb_l, tb_t, tb_w, tb_h).table
                col_ws = [Inches(2.4), Inches(0.8), Inches(1.0), Inches(1.5)]
                for ci, cw in enumerate(col_ws):
                    tb.columns[ci].width = cw
                for ci, hdr in enumerate(["問題類型", "件數", "百分比", "歸屬部門"]):
                    cell = tb.cell(0, ci)
                    cell.text = hdr
                    cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.bold  = True
                            run.font.color.rgb = WHITE
                            run.font.size  = Pt(13)
                            run.font.name  = FONT
                for ri, (_, r) in enumerate(stats.head(rows_n - 1).iterrows(), 1):
                    try:   pct = f'{int(float(r["百分比"]))}%'
                    except: pct = f'{r["百分比"]}%'
                    dept = str(r.get("歸屬部門", ""))
                    vals = [str(r["問題類型"]), str(int(r["件數"])), pct, dept]
                    # 依部門套用品牌色為列底色
                    dept_hex = DEPT_COLOR_MAP.get(dept, "")
                    if dept_hex:
                        r_bg = RGBColor(
                            int(dept_hex[1:3], 16),
                            int(dept_hex[3:5], 16),
                            int(dept_hex[5:7], 16),
                        )
                        # 淡化：混入白色 80%
                        r_bg = RGBColor(
                            min(255, int(r_bg[0] * 0.25 + 255 * 0.75)),
                            min(255, int(r_bg[1] * 0.25 + 255 * 0.75)),
                            min(255, int(r_bg[2] * 0.25 + 255 * 0.75)),
                        )
                    else:
                        r_bg = LGRAY if ri % 2 == 0 else BEIGE
                    for ci, v in enumerate(vals):
                        cell = tb.cell(ri, ci)
                        cell.text = v
                        cell.fill.solid(); cell.fill.fore_color.rgb = r_bg
                        for para in cell.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.size  = Pt(12)
                                run.font.color.rgb = DARK
                                run.font.name  = FONT

        if len(slides) >= 2:
            _fill_slide(slides[1],
                        f"{source_name} 客訴問題分析",
                        ["chart_問題類型分布.png"],
                        add_table=True)
        if len(slides) >= 3:
            _fill_slide(slides[2],
                        f"{source_name} 機台與細項分析",
                        ["chart_十大問題細項.png", "chart_機台問題占比.png"],
                        add_table=False)

    # ════════════════════════════════════════════════════════
    #  從零構建（範本不存在時）
    # ════════════════════════════════════════════════════════
    else:
        SWi = SW / 914400   # EMU → inches
        SHi = SH / 914400

        # ── Slide 1: 封面 ──
        s0 = prs.slides.add_slide(blank_layout())
        add_rect(s0, 0, 0, SWi, SHi, BLUE)      # 全藍背景
        add_text(s0, "ECOCO 客訴分析簡報",
                 1.0, SHi*0.25, SWi-2, 1.2, 36, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s0, f"報告日期：{datetime.now().strftime('%Y/%m/%d')}",
                 1.0, SHi*0.52, SWi-2, 0.5, 16,
                 color=BEIGE, align=PP_ALIGN.CENTER)
        add_text(s0, f"資料來源：{source_name}",
                 1.0, SHi*0.64, SWi-2, 0.5, 14,
                 color=BEIGE, align=PP_ALIGN.CENTER)
        add_text(s0, "凡立橙股份有限公司",
                 1.0, SHi*0.82, SWi-2, 0.4, 13,
                 color=WHITE, align=PP_ALIGN.CENTER)

        # ── Slide 2: 問題類型分析 ──
        s1 = prs.slides.add_slide(blank_layout())
        add_header(s1, f"客訴問題分析 — {source_name}",
                   f"報告日期：{datetime.now().strftime('%Y/%m/%d')}　資料來源：{source_name}")
        # 表格（左半）
        rows_n = min(len(stats) + 1, 10)
        tbl_left = Inches(0.3); tbl_top = Inches(1.15)
        tbl_w    = Inches(5.8); tbl_h   = Inches(SHi - 1.4)
        tb = s1.shapes.add_table(rows_n, 4, tbl_left, tbl_top, tbl_w, tbl_h).table
        tb.columns[0].width = Inches(2.2)
        tb.columns[1].width = Inches(0.9)
        tb.columns[2].width = Inches(1.0)
        tb.columns[3].width = Inches(1.5)
        for ci, hdr in enumerate(["問題類型", "件數", "百分比", "歸屬部門"]):
            c = tb.cell(0, ci); c.text = hdr
            c.fill.solid(); c.fill.fore_color.rgb = BLUE
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.bold = True; run.font.color.rgb = WHITE
                    run.font.size = Pt(12); run.font.name = FONT
        for ri, (_, r) in enumerate(stats.head(rows_n - 1).iterrows(), 1):
            try:   pct = f'{int(float(r["百分比"]))}%'
            except: pct = f'{r["百分比"]}%'
            vals = [str(r["問題類型"]), str(r["件數"]), pct,
                    str(r.get("歸屬部門", ""))]
            bg = LGRAY if ri % 2 == 0 else BEIGE
            for ci, v in enumerate(vals):
                c = tb.cell(ri, ci); c.text = v
                c.fill.solid(); c.fill.fore_color.rgb = bg
                for para in c.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(11); run.font.color.rgb = DARK
                        run.font.name = FONT
        # 圖表（右半）
        if chart_pack and "chart_問題類型分布.png" in chart_pack:
            add_img(s1, chart_pack["chart_問題類型分布.png"],
                    6.25, 1.15, SWi - 6.55, SHi - 1.4)

        # ── Slide 3: 機台與細項分析 ──
        s2 = prs.slides.add_slide(blank_layout())
        add_header(s2, f"機台與細項分析 — {source_name}",
                   f"報告日期：{datetime.now().strftime('%Y/%m/%d')}")
        half_w = (SWi - 0.6) / 2
        ch_t = 1.15; ch_h = SHi - 1.4
        if chart_pack and "chart_機台問題占比.png" in chart_pack:
            add_img(s2, chart_pack["chart_機台問題占比.png"],
                    0.3, ch_t, half_w, ch_h)
        if chart_pack and "chart_十大問題細項.png" in chart_pack:
            add_img(s2, chart_pack["chart_十大問題細項.png"],
                    0.3 + half_w + 0.15, ch_t, half_w, ch_h)

    # ── 最終：AI 重點分析投影片（所有路徑都加）──
    s_ai = prs.slides.add_slide(blank_layout())
    SWi2 = prs.slide_width  / 914400
    SHi2 = prs.slide_height / 914400
    # 藍色頁首
    add_rect(s_ai, 0, 0, SWi2, 1.05, BLUE)
    add_text(s_ai, "AI 重點問題分析",
             0.3, 0.08, 9.0, 0.55, 20, bold=True, color=WHITE)
    add_text(s_ai,
             f"資料來源：{source_name}　產出日期：{datetime.now().strftime('%Y/%m/%d')}",
             0.3, 0.65, 10.5, 0.35, 11, color=BEIGE)
    # 橘色左邊框裝飾
    add_rect(s_ai, 0.25, 1.15, 0.08, SHi2 - 1.35, ORANGE)
    # AI 文字框
    txb = s_ai.shapes.add_textbox(Inches(0.45), Inches(1.2),
                                   Inches(SWi2 - 0.65), Inches(SHi2 - 1.35))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for line in ai_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        is_head = line[:2] in ('1)', '2)', '3)', '4)', '5)', '一、', '二、', '三、')
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(14 if is_head else 13)
        run.font.bold  = is_head
        run.font.color.rgb = BLUE if is_head else DARK

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()



def upload_to_google_sheet(df: pd.DataFrame, credentials_json: dict, spreadsheet_id: str, worksheet_name: str) -> None:
    import gspread as _gs
    from google.oauth2.service_account import Credentials as _Creds
    # 必須同時包含 spreadsheets 和 drive scope
    scopes = [
        "https://spreadsheets.google.com/feeds",
        "https://www.googleapis.com/auth/spreadsheets",
        "https://www.googleapis.com/auth/drive",
    ]
    creds = _Creds.from_service_account_info(credentials_json, scopes=scopes)
    client = _gs.authorize(creds)
    try:
        sh = client.open_by_key(spreadsheet_id)
    except Exception as e:
        raise PermissionError(
            f"無法存取試算表（ID: {spreadsheet_id}）。\n"
            f"請確認已將試算表共用給：{credentials_json.get('client_email', '?')}\n"
            f"原始錯誤：{e}"
        )
    clean_df = _sanitize_df_for_sheet(df)
    values = [clean_df.columns.tolist()] + clean_df.values.tolist()
    try:
        ws = sh.worksheet(worksheet_name)
        ws.clear()
        ws.resize(rows=max(len(values), 100), cols=max(len(clean_df.columns), 10))
    except Exception:
        ws = sh.add_worksheet(
            title=worksheet_name,
            rows=max(len(values), 100),
            cols=max(len(clean_df.columns), 10),
        )
    if values:
        ws.update(values=values, range_name="A1")
    return ws.url if hasattr(ws, 'url') else ""


def apply_editor_changes(full_df: pd.DataFrame, edited: pd.DataFrame,
                         show_index, editor_state: dict | None = None,
                         drop_cols: list[str] | None = None) -> tuple[pd.DataFrame, list]:
    """把 data_editor 的編輯結果寫回完整資料。

    表格可能只顯示部分列（篩選、只看待複核），而 data_editor 回傳的 index
    不保證與原始資料一致，只保證「順序與顯示順序相同、新增列排在最後」。
    因此用 show_index（顯示中那些列的原始 id）依序對齊，
    並參考 widget 狀態裡的 deleted_rows 排除被刪掉的列。
    回傳 (更新後的完整資料, 實際被寫入的列 id)。
    """
    out = full_df.copy()
    if edited is None or len(edited) == 0:
        return out, []

    work = edited.drop(columns=[c for c in (drop_cols or []) if c in edited.columns],
                       errors="ignore")
    state = editor_state or {}
    deleted_pos = set(state.get("deleted_rows", []) or [])
    kept_index = [idx for pos, idx in enumerate(list(show_index)) if pos not in deleted_pos]

    n = min(len(kept_index), len(work))
    existing = work.iloc[:n].copy()
    existing.index = pd.Index(kept_index[:n])
    existing = existing.loc[existing.index.intersection(out.index)]
    if not existing.empty:
        out.update(existing)

    added = work.iloc[len(kept_index):].copy()
    if not added.empty:
        base = int(pd.to_numeric(pd.Series(out.index), errors="coerce").max() or -1) + 1
        added.index = pd.Index(range(base, base + len(added)))
        out = pd.concat([out, added])

    return out, list(existing.index)


def _toolbar_label(text: str) -> None:
    st.markdown(f"<div class='toolbar-title'>{text}</div>", unsafe_allow_html=True)


def render_editor_toolbar(df, edited, editor_row_index, marker_col, ai_col, summary) -> None:
    """表格上方的整併工具列。

    原本散在三處（表格上方的欄位管理、自訂選項，表格下方的「批次處理與儲存」，
    以及篩選列的「只看待複核」）合成一塊，減少上下捲動。
    這個函式是在 data_editor 之後才呼叫的（批次動作要讀 edited 的勾選狀態），
    但畫面位置由外層先建立的容器決定，所以仍然顯示在表格上方。
    """
    st.checkbox("只看待複核", key="editor_only_review",
                value=st.session_state.get("editor_only_review", summary["review"] > 0),
                help="只顯示需要人工確認的資料列")

    col_left, col_right = st.columns(2, gap="large")

    # ── 欄位管理 ────────────────────────────────────────────────
    with col_left:
        _toolbar_label("欄位管理")
        a1, a2 = st.columns([3, 1.2], vertical_alignment="bottom")
        new_col_name = a1.text_input("新增直立欄位", value="", key="editor_new_col",
                                     placeholder="輸入欄位名稱")
        if a2.button("新增欄位", key="editor_add_col", use_container_width=True):
            col_name = new_col_name.strip()
            if not col_name:
                st.warning("請輸入欄位名稱。")
            elif col_name in st.session_state["analysis_df"].columns:
                st.warning("欄位已存在。")
            else:
                st.session_state["analysis_df"][col_name] = ""
                st.session_state.pop("editor_table", None)
                st.rerun()

        protected_cols = {"選取", marker_col, ai_col, *META_COLUMNS}
        deletable_cols = [c for c in st.session_state["analysis_df"].columns
                          if c not in protected_cols]
        b1, b2 = st.columns([3, 1.2], vertical_alignment="bottom")
        del_col_name = b1.selectbox("選取欄位", options=deletable_cols, key="editor_delete_col")
        if b2.button("刪除整欄", key="editor_del_col", use_container_width=True):
            if del_col_name:
                st.session_state["analysis_df"] = st.session_state["analysis_df"].drop(
                    columns=[del_col_name], errors="ignore")
                st.session_state.pop("editor_table", None)
                st.rerun()

    # ── 批次問題處理 ────────────────────────────────────────────
    with col_right:
        _toolbar_label("批次問題處理")
        c1, c2 = st.columns(2)
        # accept_new_options：可從下拉選，也可以直接打字新增（combo box）
        batch_type_opts = ["(不變更)"] + combo_options(TYPE_OPTIONS, df, "問題類型", "_custom_types")
        batch_type = c1.selectbox("批次問題類型", batch_type_opts, key="batch_type_sel",
                                  accept_new_options=True, help="清單沒有的可直接輸入新增")
        valid_batch_det = ["(不變更)"]
        if batch_type != "(不變更)":
            valid_batch_det += TOPIC_DETAIL_MAP.get(batch_type, [])
        valid_batch_det += [d for d in st.session_state.get("_custom_details", [])
                            if d not in valid_batch_det]
        batch_detail = c2.selectbox("批次問題細項", valid_batch_det, key="batch_cat_sel",
                                    accept_new_options=True, help="清單沒有的可直接輸入新增")
        if batch_detail and batch_detail != "(不變更)":
            batch_detail = lower_english(batch_detail)

        d1, d2 = st.columns(2)
        apply_clicked = d1.button("套用勾選列", key="batch_apply", type="primary",
                                  use_container_width=True,
                                  help="把上面兩個下拉的設定寫進所有勾選的列")
        delete_clicked = d2.button("刪除勾選列", key="batch_delete",
                                   use_container_width=True)

    # ── 自訂選項 ────────────────────────────────────────────────
    _toolbar_label("自訂選項（下拉清單沒有的值先加進來，表格每一格就選得到）")
    e1, e2, e3 = st.columns([3, 3, 1.4], vertical_alignment="bottom")
    new_type_opt = e1.text_input("新增自訂問題類型", key="editor_new_type",
                                 placeholder="下拉清單沒有時輸入新增")
    new_detail_opt = e2.text_input("新增自訂問題細項", key="editor_new_detail",
                                   placeholder="英文會自動轉小寫")
    if e3.button("加入選項", key="editor_add_option", use_container_width=True):
        added = []
        if new_type_opt.strip():
            st.session_state.setdefault("_custom_types", [])
            name = new_type_opt.strip()
            if name not in st.session_state["_custom_types"]:
                st.session_state["_custom_types"].append(name)
                added.append(name)
        if new_detail_opt.strip():
            st.session_state.setdefault("_custom_details", [])
            name = lower_english(new_detail_opt.strip())
            if name not in st.session_state["_custom_details"]:
                st.session_state["_custom_details"].append(name)
                added.append(name)
        if added:
            st.session_state.pop("editor_table", None)
            st.rerun()
        else:
            st.warning("請先輸入要新增的選項名稱。")

    # ── 批次動作 ────────────────────────────────────────────────
    has_selection = "選取" in edited.columns and bool(edited["選取"].any())

    if apply_clicked:
        if not has_selection:
            st.warning("請先在表格內勾選要處理的資料列！")
        else:
            mask = edited["選取"] == True
            if batch_type != "(不變更)":
                edited.loc[mask, "問題類型"] = batch_type
                edited.loc[mask, "部門"] = edited.loc[mask, "問題類型"].map(DEPT_MAP).fillna("")
            if batch_detail != "(不變更)":
                edited.loc[mask, "問題細項"] = batch_detail
            # 細項與類型對不上時自動修正；使用者自行新增的細項不動它
            custom_details = set(st.session_state.get("_custom_details", []))

            def _fix_detail(r):
                detail = r["問題細項"]
                allowed = TOPIC_DETAIL_MAP.get(r["問題類型"], [])
                if detail in allowed or detail in custom_details:
                    return detail
                return allowed[0] if allowed else detail

            edited["問題細項"] = edited.apply(_fix_detail, axis=1)
            merged, _ = apply_editor_changes(
                st.session_state["analysis_df"], edited, editor_row_index,
                editor_state=st.session_state.get("editor_table"),
                drop_cols=[marker_col],
            )
            st.session_state["analysis_df"] = merged
            st.session_state.pop("editor_table", None)
            st.session_state["_batch_applied"] = True
            st.rerun()

    if st.session_state.pop("_batch_applied", False):
        st.success("已套用批次編輯。")

    if delete_clicked:
        if not has_selection:
            st.warning("請先在表格內勾選要刪除的資料列！")
        else:
            drop_ids = [rid for rid, flag in zip(editor_row_index, edited["選取"]) if bool(flag)]
            st.session_state["analysis_df"] = (
                st.session_state["analysis_df"].drop(index=drop_ids, errors="ignore").copy()
            )
            st.session_state.pop("editor_table", None)
            st.success("已刪除勾選列。")
            st.rerun()


def mask_phone_columns(df: pd.DataFrame) -> pd.DataFrame:
    """遮蔽「帳號手機」這類整格號碼欄位。

    新分析的資料在 core.analyze_dataframe 就已遮蔽，
    這裡是給功能三／功能四讀進來的舊歷史紀錄補上。
    """
    if df is None or df.empty:
        return df
    out = df.copy()
    for col in out.columns:
        if any(k in str(col).lower() for k in PHONE_COL_HINTS):
            out[col] = out[col].map(mask_phone_value)
    return out


def combo_options(base: list[str], df: pd.DataFrame | None, column: str,
                  session_key: str) -> list[str]:
    """下拉選項＝內建清單＋資料實際值＋使用者自行輸入過的值。

    st.column_config.SelectboxColumn 不支援直接在格子裡打字，
    所以自訂值改由表格上方的輸入框加入，加進來之後每一格都選得到，
    效果等同可打字的 combo box。
    """
    options = options_with_data_values(base, df, column)
    for name in st.session_state.get(session_key, []):
        if name and name not in options:
            options.append(name)
    return options


def options_with_data_values(base: list[str], df: pd.DataFrame | None, column: str) -> list[str]:
    """下拉選項＝內建清單＋資料裡實際出現的值。

    分類法調整（新增細項、廢除類型）後，舊資料仍可能帶著舊標記；
    若選項少了實際值，data_editor 那一格會顯示不出來甚至報錯。
    """
    options = list(base)
    if df is not None and column in df.columns:
        for v in df[column].dropna().unique():
            name = str(v).strip()
            if name and name not in options:
                options.append(name)
    return options


def dept_options_for(df: pd.DataFrame | None = None) -> list[str]:
    """部門選項＝內建清單＋知識庫學到的＋資料裡實際出現的。"""
    options = list(DEPT_OPTIONS)
    kb = get_knowledge(st.session_state.get("_knowledge_version", ""))
    learned = list(getattr(kb, "dept_by_topic", {}).values()) if kb else []
    in_data = []
    if df is not None and "部門" in df.columns:
        in_data = [str(v).strip() for v in df["部門"].dropna().unique() if str(v).strip()]
    for name in learned + in_data:
        if name and name not in options:
            options.append(name)
    return options


def render_knowledge_panel() -> None:
    """顯示分類知識庫狀態（從歷史紀錄學到什麼），並提供重建按鈕。"""
    kb = get_knowledge(st.session_state.get("_knowledge_version", ""))
    if kb is None:
        with st.expander("🧠 分類知識庫（尚未建立）", expanded=False):
            st.caption("目前只使用內建關鍵字規則。等歷史紀錄累積後，"
                       "系統會自動從過往（尤其是人工修正過的）標記學出規則與相似案例池。")
            msg = st.session_state.get("_knowledge_error")
            if msg:
                st.caption(f"讀取歷史時的訊息：{msg}")
        return

    s = kb.stats
    with st.expander(
        f"🧠 分類知識庫：已從 {s.get('history_rows', 0)} 筆歷史標記學到 "
        f"{s.get('rules', 0)} 條規則、{s.get('fingerprints', 0)} 組指紋",
        expanded=False,
    ):
        k1, k2, k3, k4 = st.columns(4)
        k1.metric("歷史標記", s.get("history_rows", 0))
        k2.metric("人工確認", s.get("confirmed_rows", 0))
        k3.metric("自動挖掘規則", s.get("rules", 0))
        k4.metric("相似案例池", s.get("knn_examples", 0))
        if kb.rules:
            st.dataframe(
                pd.DataFrame([
                    {"問題類型": r.topic, "問題細項": r.detail,
                     "關鍵字": "、".join(r.terms[:6]), "歷史筆數": r.support,
                     "實測準確率": f"{r.precision:.0%}"}
                    for r in kb.rules[:15]
                ]),
                use_container_width=True, hide_index=True,
            )
        if kb.dept_by_topic:
            st.caption("從歷史學到的部門對應：" +
                       "　".join(f"{t}→{d}" for t, d in kb.dept_by_topic.items()))
            conflicts = {t: d for t, d in kb.dept_by_topic.items()
                         if DEPT_MAP.get(t) and DEPT_MAP.get(t) != d}
            if conflicts:
                st.warning(
                    "歷史實際填的部門與程式內建的 DEPT_MAP 不一致：" +
                    "　".join(f"{t}：內建「{DEPT_MAP[t]}」／歷史「{d}」" for t, d in conflicts.items()) +
                    "。目前以內建為準；若公司已改部門編制，把環境變數 "
                    "PREFER_LEARNED_DEPT 設為 true 即可改以歷史名稱為準。"
                )
        if st.button("重建知識庫", key="kb_rebuild"):
            get_knowledge.clear()
            bump_knowledge_version()
            st.rerun()


def section_1():
    page_header("功能一：檔案上傳與分析區",
                "支援上傳 excel / csv / pdf，分析並產出【問題類型、問題細項】。")
    render_knowledge_panel()

    # File info badge — no long text, just a compact pill with truncated name
    if st.session_state.get("_uploaded_bytes") and st.session_state.get("_uploaded_name"):
        fname_short = st.session_state['_uploaded_name']
        if len(fname_short) > 30:
            fname_short = fname_short[:14] + "..." + fname_short[-12:]
        st.markdown(
            f"<span class='file-badge'>&#128196; {fname_short}</span>",
            unsafe_allow_html=True
        )

    uploaded = st.file_uploader("上傳新檔案", type=["xlsx", "xls", "csv", "pdf"], key="uploader")
    # Persist file bytes across menu switches
    if uploaded is not None:
        if uploaded.name != st.session_state.get("_uploaded_name"):
            st.session_state.pop("_editing_history_id", None)
            st.session_state.pop("_saved_history_id", None)
        st.session_state["_uploaded_bytes"] = uploaded.read()
        st.session_state["_uploaded_name"] = uploaded.name
        st.session_state["_uploaded_type"] = uploaded.type

    # Restore from session if user switched tabs and came back
    if uploaded is None and st.session_state.get("_uploaded_bytes") is not None:
        saved_name = st.session_state.get("_uploaded_name", "file")
        buf = io.BytesIO(st.session_state["_uploaded_bytes"])
        df_raw_bytes = load_input_file(buf, filename=saved_name)
        st.caption(f"已載入 {saved_name}（從記憶復原），資料筆數：{len(df_raw_bytes)}")
        df_raw = make_unique_columns(df_raw_bytes)
        uploaded_name = saved_name
    elif uploaded is not None:
        fname = st.session_state.get("_uploaded_name", uploaded.name)
        df_raw = make_unique_columns(load_input_file(
            io.BytesIO(st.session_state["_uploaded_bytes"]), filename=fname
        ))
        uploaded_name = uploaded.name
        st.caption(f"已載入 {uploaded.name}，資料筆數：{len(df_raw)}")
    else:
        if "analysis_df" not in st.session_state:
            st.info("請上傳檔案開始分析。")
            return
        # Already analysed, show results without needing the raw file
        df_raw = None
        uploaded_name = st.session_state.get("source_name", "")

    if df_raw is not None:
        cols = list(df_raw.columns)
        if not cols:
            st.warning("檔案沒有可用欄位。")
            return

        # ── 欄位自動偵測（偵測得到就不必人工指定）──
        det = detect_columns(df_raw)
        if det.ok:
            st.markdown(
                "<div class='ecoco-card' style='border-left:4px solid #060E9F;'>"
                f"✅ 已自動判斷欄位　主題：<b>{det.subject}</b>　內容：<b>{det.content}</b>"
                f"　日期：<b>{det.date or '（無）'}</b></div>",
                unsafe_allow_html=True,
            )
        else:
            st.warning("無法自動判斷欄位對應，請在下方手動指定。")

        def _idx(name, options, fallback=0):
            return options.index(name) if name in options else fallback

        # 日期欄位不再讓人工指定，一律採用自動偵測結果（偵測不到就不帶日期）
        date_col = det.date if det.date in cols else "(無)"
        with st.expander("欄位對應（自動判斷結果，可手動調整）", expanded=not det.ok):
            subject_col = st.selectbox("用戶填寫的主題欄位", options=cols,
                                       index=_idx(det.subject, cols, 0), key="col_subject")
            content_col = st.selectbox("用戶內容欄位", options=cols,
                                       index=_idx(det.content, cols, min(1, len(cols) - 1)),
                                       key="col_content")
            for field_name, label in (("subject", "主題"), ("content", "內容"), ("date", "日期")):
                why = det.reasons.get(field_name)
                if why:
                    st.caption(f"{label}欄判斷依據：{why}"
                               f"（信心 {det.confidence.get(field_name, 0):.0%}）")

        cfg = AnalysisConfig(subject_col=subject_col, content_col=content_col,
                             date_col=None if date_col == "(無)" else date_col)

        # ── 自動分析：同一份檔案與同一組設定只會跑一次 ──
        run_sig = f"{uploaded_name}|{len(df_raw)}|{subject_col}|{content_col}|{date_col}"
        manual_run = st.button("重新分析", help="重新套用目前設定與最新分類規則")
        auto_run = (auto_config.auto_analyze_on_upload() and det.ok
                    and st.session_state.get("_auto_run_sig") != run_sig)
        if manual_run or auto_run:
            work = df_raw.copy()
            with st.spinner("正在自動分類，請稍候…"):
                result = analyze_dataframe(work, cfg)
            st.session_state["analysis_df"] = result
            st.session_state["source_name"] = uploaded_name
            st.session_state["_auto_run_sig"] = run_sig
            st.session_state.pop("editor_table", None)

            summary = review_summary(result)
            st.success(
                f"已自動分析 {summary['total']} 筆：完全自動採用 {summary['auto']} 筆"
                f"（{summary['auto_rate']:.0%}），待人工複核 {summary['review']} 筆"
            )
            if auto_config.auto_save_history():
                existing_id = (st.session_state.get("_editing_history_id")
                               or st.session_state.get("_saved_history_id", ""))
                try:
                    _, _, history_id = save_history(result, uploaded_name, existing_id=existing_id)
                    st.session_state["_saved_history_id"] = history_id
                    if st.session_state.get("_gsheet_error"):
                        st.warning(st.session_state["_gsheet_error"])
                    else:
                        st.caption(f"已自動存入歷史紀錄（{history_id}）")
                        bump_knowledge_version()
                except Exception as exc:
                    st.warning(f"自動存檔失敗，可稍後手動儲存：{str(exc)[:200]}")

    if "analysis_df" not in st.session_state:
        return
    df = st.session_state["analysis_df"]
    # ── 自動化與審核總覽 ──
    summary = review_summary(df)
    causes = summary.get("review_causes", {})
    n_detail_only = causes.get("僅需確認細項", 0)
    n_full = causes.get("信心不足", 0) + causes.get("各層判斷分歧", 0)
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("總筆數", summary["total"])
    m2.metric("完全自動採用", f"{summary['auto']} 筆", f"{summary['auto_rate']:.0%}")
    m3.metric("僅需確認細項", f"{n_detail_only} 筆",
              help="類型與部門系統有把握，只有細項要人工挑一個")
    m4.metric("需完整判斷", f"{n_full} 筆", help="類型也沒把握或各層分歧")
    if causes.get("稽核抽樣"):
        st.caption(f"🔍 另有 {causes['稽核抽樣']} 筆是從自動採用中抽出的品質抽驗")
    if summary.get("agreement"):
        st.caption("交叉驗證：" + "　".join(f"{k} {v} 筆" for k, v in summary["agreement"].items()))

    # 篩選條件已整併到表格上方的工具列（只保留「只看待複核」）。
    only_review = bool(st.session_state.get("editor_only_review", summary["review"] > 0))

    # 不重設 index：表格顯示由 hide_index 負責，
    # 原始 index 必須保留，儲存時才知道每一列要寫回哪裡。
    show = make_unique_columns(df.copy())
    if only_review and "_needs_review" in show.columns:
        show = show[show["_needs_review"].fillna(False).astype(bool)]
        # 排序：需完整判斷 → 只需挑細項 → 稽核抽樣；同組內最沒把握的優先
        cause_order = {"信心不足": 0, "各層判斷分歧": 0, "僅需確認細項": 1, "稽核抽樣": 2}
        if "_review_cause" in show.columns:
            ordered = show.assign(_order=show["_review_cause"].map(cause_order).fillna(0))
            show = ordered.sort_values(["_order", "_confidence"], kind="stable").drop(columns=["_order"])

    st.markdown('<div class="editor-toolbar-title">可編輯標記表（支援下拉 + 手動編輯）</div>', unsafe_allow_html=True)

    # ---- 待複核標示 ---
    ai_col = "_ai_filled"
    MARKER_COL = "AI標記"  # kept for save compatibility only
    has_ai_col = ai_col in show.columns
    review_col = "_needs_review"
    has_review_col = review_col in show.columns
    n_review = int(show[review_col].fillna(False).astype(bool).sum()) if has_review_col else 0
    n_ai = int(show[ai_col].fillna(False).astype(bool).sum()) if has_ai_col else 0

    if has_review_col:
        if n_review:
            st.markdown(
                f"""
                <div style='background:#fff5f5; border:1px solid #ffb3b3; border-radius:8px;
                            padding:8px 14px; margin-bottom:8px; font-size:0.85rem;'>
                  <b style='color:#cc0000;'>⚠ 待人工複核 {n_review} 筆</b>（已按信心由低到高排序）。三種原因：
                  <b>信心不足</b>＝各層都沒把握；<b>各層判斷分歧</b>＝規則與相似案例給出不同答案
                  （實測這類的細項正確率只有三成，務必看）；<b>🔍 稽核抽樣</b>＝系統其實有把握，
                  抽出來抽驗品質用。其餘資料已自動採用，不需逐列檢查。
                </div>
                """,
                unsafe_allow_html=True,
            )
        else:
            st.success("✅ 顯示中的資料都不需要人工複核。")
    elif n_ai > 0:
        st.markdown(
            f"""
            <div style='background:#fff5f5; border:1px solid #ffb3b3; border-radius:8px;
                        padding:8px 14px; margin-bottom:8px; font-size:0.85rem;'>
              <b style='color:#cc0000;'>● AI 自動標記</b>：共 <b style='color:#cc0000;'>{n_ai} 筆</b> 原始欄位空白或無效，
              已由 AI 根據客訴內容自動分析填入。
              請針對這幾筆核對，如需修改請直接在表格中下拉選擇，再點「💾 儲存修改」確認。
            </div>
            """,
            unsafe_allow_html=True
        )

    st.caption("💡 直接在表格中下拉選擇問題類型 / 問題細項，調整完成後點擊「💾 儲存修改」。")

    # 工具列（檢視 / 欄位管理 / 批次問題處理 / 自訂選項）先佔位，
    # 實際內容在 data_editor 之後才填入 —— 批次動作需要 edited 的勾選狀態，
    # 但版面上它要出現在表格上方。
    toolbar = st.container(key="editor_toolbar")

    # 顯示用欄位：隱藏所有底線開頭的內部欄位（信心、判斷來源等稽核資訊）
    display_cols = [c for c in visible_columns(show) if c not in (MARKER_COL, "選取")]
    editor_row_index = list(show.index)   # 顯示中每一列對應的原始資料列 id
    show_display = show[display_cols].reset_index(drop=True)

    # 備註欄：顯示複核原因與把握度，讓人知道為什麼這列要看
    if has_review_col:
        def _marker(row):
            if not bool(row.get(review_col, False)):
                return ""
            conf = float(row.get("_confidence", 0) or 0)
            tconf = float(row.get("_topic_confidence", 0) or 0)
            cause = str(row.get("_review_cause", "") or "待複核")
            if cause == "稽核抽樣":
                return f"🔍 稽核抽樣（信心 {conf:.0%}）"
            if cause == "僅需確認細項":
                return f"✏ 只需挑細項（類型把握 {tconf:.0%}，細項把握 {conf:.0%}）"
            agree = str(row.get("_agreement", "") or "")
            extra = f"／{agree}" if agree else ""
            return f"⚠ {cause}（信心 {conf:.0%}{extra}）"
        marker_vals = list(show.apply(_marker, axis=1)) if len(show) else []
    elif has_ai_col:
        marker_vals = list(show[ai_col].map(lambda x: "⭐(AI填寫)" if x else ""))
    else:
        marker_vals = [""] * len(show_display)
        
    # 「選取」固定在第一欄，「備註」緊接其後。
    # 原本沒有這一欄，批次套用／刪除勾選列因此永遠抓不到勾選狀態。
    prev_sel = show["選取"] if "選取" in show.columns else None
    sel_vals = list(prev_sel.fillna(False).astype(bool)) if prev_sel is not None else [False] * len(show_display)
    show_display.insert(0, "選取", sel_vals)
    show_display.insert(1, MARKER_COL, marker_vals)

    edited = st.data_editor(
        show_display,
        use_container_width=True,
        num_rows="fixed",
        hide_index=True,
        column_config={
            "選取": st.column_config.CheckboxColumn("選取", help="勾選要批次處理的列", pinned=True),
            MARKER_COL: st.column_config.TextColumn("備註", disabled=True),
            "問題類型": st.column_config.SelectboxColumn(
                options=combo_options(TYPE_OPTIONS, df, "問題類型", "_custom_types"), required=True),
            "問題細項": st.column_config.SelectboxColumn(
                options=combo_options(DETAIL_OPTIONS, df, "問題細項", "_custom_details"), required=True),
            "部門": st.column_config.SelectboxColumn(options=dept_options_for(df)),
        },
        key="editor_table",
    )

    with toolbar:
        render_editor_toolbar(df, edited, editor_row_index, MARKER_COL, ai_col, summary)

    # 儲存按鈕在表格下方
    sv_col1, sv_col2, sv_col3 = st.columns([2, 2, 6])
    if sv_col2.button("✔ 全部接受系統判斷", use_container_width=True,
                      help="把目前顯示的列一次標記為已複核（不改內容），適合快速掠過沒問題的列"):
        accepted = st.session_state["analysis_df"].copy()
        idx = [i for i in editor_row_index if i in accepted.index]
        for col, val in (("_needs_review", False), ("_confidence", 1.0),
                         ("_source_layer", "人工確認"), ("_review_cause", "")):
            if col in accepted.columns:
                accepted.loc[idx, col] = val
        if "_ai_filled" in accepted.columns:
            accepted.loc[idx, "_ai_filled"] = False
        st.session_state["analysis_df"] = accepted
        st.session_state.pop("editor_table", None)
        src_name = st.session_state.get("source_name", "未命名")
        existing_id = (st.session_state.get("_editing_history_id")
                       or st.session_state.get("_saved_history_id", ""))
        try:
            _, _, history_id = save_history(accepted, src_name, existing_id=existing_id)
            st.session_state["_saved_history_id"] = history_id
            bump_knowledge_version()
        except Exception as exc:
            st.warning(f"存檔失敗：{str(exc)[:200]}")
        st.session_state["_accepted_n"] = len(idx)
        st.rerun()

    if st.session_state.pop("_accepted_n", 0):
        st.success("已接受系統判斷並存檔，這些列已納入知識庫")

    if sv_col1.button("💾 儲存修改", use_container_width=True):
        full_df, touched = apply_editor_changes(
            st.session_state["analysis_df"], edited, editor_row_index,
            editor_state=st.session_state.get("editor_table"),
            drop_cols=["選取", MARKER_COL],
        )
        # 人工已看過這些列 → 解除待複核，記為人工確認（下次分類的黃金標註）
        for col, val in (("_needs_review", False), ("_confidence", 1.0),
                         ("_source_layer", "人工確認"), ("_review_cause", "")):
            if col in full_df.columns:
                full_df.loc[touched, col] = val
        if "_ai_filled" in full_df.columns:
            full_df["_ai_filled"] = False
        st.session_state["analysis_df"] = full_df
        # Also push to drafts list
        src_name = st.session_state.get("source_name", "未命名")
        if "_draft_list" not in st.session_state:
            st.session_state["_draft_list"] = []
        # Avoid duplicate same name drafts – update existing
        draft_ids = [d["name"] for d in st.session_state["_draft_list"]]
        if src_name not in draft_ids:
            st.session_state["_draft_list"].insert(0, {"name": src_name, "df": full_df.copy()})
        else:
            for d in st.session_state["_draft_list"]:
                if d["name"] == src_name:
                    d["df"] = full_df.copy()
        existing_id = st.session_state.get("_editing_history_id") or st.session_state.get("_saved_history_id", "")
        _, _, history_id = save_history(full_df, src_name, existing_id=existing_id)
        st.session_state["_saved_history_id"] = history_id
        if st.session_state.get("_gsheet_error"):
            st.warning(st.session_state["_gsheet_error"])
        else:
            bump_knowledge_version()   # 人工修正納入知識庫，下次分類更準
            st.success(f"已儲存「{src_name}」，並已納入分類知識庫")

    # 已儲存草稿列表
    if st.session_state.get("_draft_list"):
        st.markdown("---")
        st.markdown("##### 已儲存的草稿")
        for idx, draft in enumerate(st.session_state["_draft_list"]):
            d_col1, d_col2, d_col3, d_col4 = st.columns([5, 1, 1, 1])
            d_col1.markdown(
                f"<div style='padding-top:0.45rem; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; font-weight:600;'>"
                f"📄 {draft['name']}</div>",
                unsafe_allow_html=True
            )
            if d_col2.button("[載入]", key=f"draft_load_{idx}", use_container_width=True):
                st.session_state["analysis_df"] = draft["df"].copy()
                st.session_state["source_name"] = draft["name"]
                st.success(f"已載入「{draft['name']}」，可繼續編輯。")
            if d_col3.button("[修改]", key=f"draft_edit_{idx}", use_container_width=True):
                st.session_state["analysis_df"] = draft["df"].copy()
                st.session_state["source_name"] = draft["name"]
                st.rerun()
            if d_col4.button("[X]", key=f"draft_del_{idx}", use_container_width=True):
                st.session_state["_draft_list"].pop(idx)
                st.rerun()

    final_df = st.session_state["analysis_df"]
    
    st.markdown("#### 下載分析結果 (下載後自動歸檔至歷史紀錄)")
    dl_format = st.radio("選擇下載格式", ["Excel", "CSV", "PDF"], horizontal=True)
    
    def on_download():
        existing_id = st.session_state.pop("_editing_history_id", "") or st.session_state.get("_saved_history_id", "")
        _, _, history_id = save_history(final_df, st.session_state.get("source_name", "unknown"), existing_id=existing_id)
        st.session_state["_saved_history_id"] = history_id
        st.session_state["history_saved_msg"] = True

    if dl_format == "Excel":
        out_name = f"{datetime.now().strftime('%Y%m%d')}_分析.xlsx"
        data_bytes = to_excel_bytes(final_df)
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif dl_format == "CSV":
        out_name = f"{datetime.now().strftime('%Y%m%d')}_分析.csv"
        data_bytes = to_csv_bytes(final_df)
        mime = "text/csv"
    else:
        out_name = f"{datetime.now().strftime('%Y%m%d')}_分析單.pdf"
        try:
            dl_key = f"_pdf_download_count_{st.session_state.get('source_name', 'unknown')}"
            st.session_state[dl_key] = int(st.session_state.get(dl_key, 0)) + 1
            data_bytes = to_pdf_bytes(final_df, st.session_state.get("source_name", "unknown"), st.session_state[dl_key])
            mime = "application/pdf"
        except Exception as e:
            st.error(f"PDF 產生錯誤: {e}")
            data_bytes = b""
            mime = "application/pdf"

    st.download_button(
        label=f"📥 下載 {dl_format} 格式分析",
        data=data_bytes,
        file_name=out_name,
        mime=mime,
        on_click=on_download
    )
    
    if st.session_state.get("history_saved_msg"):
        if st.session_state.get("_gsheet_error"):
            st.warning(st.session_state["_gsheet_error"])
        else:
            st.success("檔案已下載，並自動保存至歷史紀錄。")
        st.session_state["history_saved_msg"] = False

    st.markdown("#### 分析文字產出")
    summary_text = generate_ai_summary(final_df)
    st.text_area("分析結果文字", summary_text, height=120)
    st.download_button(
        "下載分析文字（txt）",
        data=summary_text.encode("utf-8"),
        file_name=f"{datetime.now().strftime('%Y%m%d')}_分析文字.txt",
        mime="text/plain",
    )

    with st.expander("上傳到 Google Sheet"):
        st.write("請提供 Service Account JSON 與 Spreadsheet ID")
        cred_file = st.file_uploader("Google Service Account JSON", type=["json"], key="gcp_json")
        spreadsheet_id = st.text_input("Spreadsheet ID")
        ws_name = st.text_input("Worksheet 名稱", value=datetime.now().strftime("%Y%m%d_分析"))
        if st.button("上傳 Google Sheet"):
            if not cred_file or not spreadsheet_id:
                st.error("請先上傳 JSON 並填寫 Spreadsheet ID。")
            else:
                try:
                    credentials_json = json.loads(cred_file.getvalue().decode("utf-8"))
                    upload_to_google_sheet(final_df, credentials_json, spreadsheet_id, ws_name)
                    st.success(f"✅ 已上傳到 Google Sheet 工作表：{ws_name}")
                    st.info(f"📋 Service Account：{credentials_json.get('client_email', '')}")
                except PermissionError as e:
                    st.error(str(e))
                    st.warning("👉 請到 Google 試算表右上角「共用」，加入上方 Service Account email，並給予「編輯者」權限")
                except Exception as e:
                    st.error(f"上傳失敗：{e}")


def render_charts_from_stats(stats: pd.DataFrame, df: pd.DataFrame, key_prefix: str = ""):
    """Render interactive Plotly charts with per-chart color pickers."""

    # ── 顏色設定 expander ──────────────────────────────────────────
    kp = key_prefix or "main"
    with st.expander("🎨 調整圖表顏色（可個別修改）", expanded=False):
        ca, cb, cc = st.columns(3)
        # 問題類型直條圖：預設「依部門品牌色」，勾選後可指定單色
        use_single_bar = ca.checkbox("直條圖使用單一顏色", key=f"{kp}_cb_bar")
        c_bar_single   = ca.color_picker("直條圖顏色", value=BRAND_ORANGE, key=f"{kp}_cp_bar") if use_single_bar else None

        # 圓餅圖：最多3個扇形獨立調色
        pie_c1 = cb.color_picker("圓餅圖 第1色（主）", value=BRAND_BLUE,   key=f"{kp}_cp_pie1")
        pie_c2 = cb.color_picker("圓餅圖 第2色（次）", value=BRAND_ORANGE, key=f"{kp}_cp_pie2")
        pie_c3 = cb.color_picker("圓餅圖 第3色",       value=BRAND_LBLUE,  key=f"{kp}_cp_pie3")

        c_hbar = cc.color_picker("細項橫條圖顏色", value=BRAND_BLUE, key=f"{kp}_cp_hbar")

    custom_pie   = [pie_c1, pie_c2, pie_c3] + BRAND_PALETTE[3:]
    custom_hbar  = c_hbar

    # ── 圓餅圖資料（供 Plotly + matplotlib 共用）────────────────
    df_machine = df[df["問題類型"] == "機台問題類型"].copy()
    m_stats = None
    if not df_machine.empty:
        def _gmt(row):
            txt = str(row.get("用戶內容", "")) + " " + str(row.get("主旨", ""))
            if "方舟" in txt: return "方舟站"
            if "電池" in txt: return "電池機"
            return "收瓶機"
        df_machine["機台機型"] = df_machine.apply(_gmt, axis=1)
        m_stats = df_machine["機台機型"].value_counts().reset_index()
        m_stats.columns = ["機型", "件數"]

    detail_stats = df["問題細項"].value_counts().reset_index().head(10)
    detail_stats.columns = ["問題細項", "件數"]

    c1, c2, c3 = st.columns(3)

    # ── 圖1：問題類型直條圖 ────────────────────────────────────
    if use_single_bar:
        fig1 = px.bar(stats, x="問題類型", y="件數", text="百分比",
                      title="問題類型分布", color_discrete_sequence=[c_bar_single])
        fig1.update_traces(marker_color=c_bar_single)
    else:
        fig1 = px.bar(stats, x="問題類型", y="件數",
                      color="歸屬部門", text="百分比", title="問題類型分布",
                      color_discrete_map=DEPT_COLOR_MAP)
    fig1.update_traces(texttemplate="%{text}%", textposition="outside")
    if not use_single_bar:
        add_counts_to_legend(fig1, stats.groupby("歸屬部門")["件數"].sum())
    fig1.update_layout(height=420, yaxis=dict(tickformat="d", nticks=6),
                       legend=BAR_LEGEND, margin=dict(t=45, b=0, r=110))
    c1.plotly_chart(fig1, use_container_width=True, key=f"{kp}_fig1")

    # ── 圖2：機台圓餅圖 ────────────────────────────────────────
    if m_stats is not None:
        m_counts = m_stats.set_index("機型")["件數"]
        m_labels = pie_legend_labels(m_counts)
        cmap = {label: custom_pie[i % len(custom_pie)] for i, label in enumerate(m_labels)}
        fig2 = px.pie(names=m_labels, values=list(m_counts.values),
                      title="機台問題細分比較", hole=0.3,
                      color=m_labels, color_discrete_map=cmap)
        fig2.update_traces(texttemplate="%{percent:.1%}", textinfo="percent")
        fig2.update_layout(height=420, showlegend=True,
                           legend=PIE_LEGEND, margin=dict(t=45, b=0, l=0, r=130))
        c2.plotly_chart(fig2, use_container_width=True, key=f"{kp}_fig2")
    else:
        with c2:
            empty_state("沒有資料紀錄　—　無機台相關數據")

    # ── 圖3：十大細項橫條圖 ────────────────────────────────────
    fig3 = px.bar(detail_stats, x="件數", y="問題細項",
                  orientation="h", title="十大問題細項分布",
                  color_discrete_sequence=[custom_hbar])
    fig3.update_traces(marker_color=custom_hbar)
    fig3.update_layout(height=420, yaxis={"categoryorder": "total ascending"},
                       xaxis=dict(tickformat="d", nticks=6),
                       margin=dict(t=45, b=0, l=0, r=0))
    c3.plotly_chart(fig3, use_container_width=True, key=f"{kp}_fig3")

    # ── 把用戶自選顏色存進 session_state 供 PPT/ZIP 使用 ────────
    st.session_state[f"chart_colors_{kp}"] = {
        "bar":  c_bar_single if use_single_bar else None,
        "pie":  custom_pie,
        "hbar": custom_hbar,
    }


def render_charts(df: pd.DataFrame, key_prefix: str = ""):
    date_cols = [c for c in df.columns if "日期" in c or "date" in c.lower()]
    if date_cols:
        dcol = date_cols[0]
        try:
            df[dcol] = pd.to_datetime(df[dcol], errors="coerce")
            valid_dates = df[dcol].dropna()
            if not valid_dates.empty:
                min_d = valid_dates.min().date()
                max_d = valid_dates.max().date()
                st.markdown("##### 分析日期區間")
                c_d1, c_d2 = st.columns(2)
                start_d = c_d1.date_input("起始日期", value=min_d, min_value=min_d, max_value=max_d, key=f"{key_prefix}_sd")
                end_d   = c_d2.date_input("結束日期", value=max_d, min_value=min_d, max_value=max_d, key=f"{key_prefix}_ed")
                df = df[(df[dcol].dt.date >= start_d) & (df[dcol].dt.date <= end_d)]
        except Exception:
            pass

    stats = df["問題類型"].value_counts().rename_axis("問題類型").reset_index(name="件數")
    stats["百分比"] = (stats["件數"] / max(stats["件數"].sum(), 1) * 100).round(0).astype(int)
    stats["歸屬部門"] = stats["問題類型"].map(DEPT_MAP).fillna("未分配")

    c1, c2, c3 = st.columns(3)
    
    fig1 = px.bar(
        stats, x="問題類型", y="件數", color="歸屬部門", text="百分比", title="問題類型分布",
        color_discrete_sequence=["#FF5000", "#060E9F", "#FFCE00", "#8EB9C9", "#0076A9", "#FAE0B8"]
    )
    fig1.update_traces(texttemplate="%{text}%", textposition="outside")
    add_counts_to_legend(fig1, stats.groupby("歸屬部門")["件數"].sum())
    fig1.update_layout(height=400, legend=BAR_LEGEND, margin=dict(t=45, b=0, r=110))
    c1.plotly_chart(fig1, use_container_width=True, key=f"{key_prefix}_fig1" if key_prefix else None)

    df_machine = df[df["問題類型"] == "機台問題類型"].copy()
    if not df_machine.empty:
        def get_machine_type(row):
            txt = str(row.get("用戶內容", "")) + " " + str(row.get("主旨", ""))
            if "方舟" in txt: return "方舟站"
            if "電池" in txt: return "電池機"
            return "收瓶機"
        df_machine["機台機型"] = df_machine.apply(get_machine_type, axis=1)
        m_stats = df_machine["機台機型"].value_counts().reset_index()
        m_stats.columns = ["機型", "件數"]
        m_counts = m_stats.set_index("機型")["件數"]
        m_labels = pie_legend_labels(m_counts)
        color_map = {label: BRAND_PALETTE[i % len(BRAND_PALETTE)]
                     for i, label in enumerate(m_labels)}
        fig2 = px.pie(
            names=m_labels, values=list(m_counts.values),
            title="機台問題細分比較", hole=0.3,
            color=m_labels, color_discrete_map=color_map,
        )
        fig2.update_traces(texttemplate="%{percent:.1%}", textinfo="percent")
        fig2.update_layout(height=400, showlegend=True,
                           legend=PIE_LEGEND, margin=dict(t=40, b=0, l=0, r=130))
        c2.plotly_chart(fig2, use_container_width=True, key=f"{key_prefix}_fig2" if key_prefix else None)
    else:
        with c2:
            empty_state("沒有資料紀錄　—　無機台相關數據")

    detail_stats = df["問題細項"].value_counts().reset_index().head(10)
    detail_stats.columns = ["問題細項", "件數"]
    fig3 = px.bar(
        detail_stats, x="件數", y="問題細項",
        orientation="h", title="十大問題細項分布",
        color_discrete_sequence=[BRAND_BLUE],
    )
    fig3.update_traces(marker_color=BRAND_BLUE)
    fig3.update_layout(
        height=400,
        yaxis={"categoryorder": "total ascending"},
        xaxis=dict(tickformat="d", nticks=6),
        margin=dict(t=40, b=0, l=0, r=0),
    )
    c3.plotly_chart(fig3, use_container_width=True, key=f"{key_prefix}_fig3" if key_prefix else None)


def section_2():
    page_header("功能二：圖表化與 AI 重點分析",
                "各問題類型件數與百分比、歸屬部門，可預覽與下載 AI 重點分析。")
    if "analysis_df" not in st.session_state:
        st.info("請先在功能一完成分析。")
        return
    df_full = st.session_state["analysis_df"]
    if df_full.empty:
        st.warning("目前沒有資料。")
        return

    # --- Date range filter ---
    date_cols = [c for c in df_full.columns if "日期" in c or "date" in c.lower()]
    df = df_full.copy()
    start_d = end_d = None
    if date_cols:
        dcol = date_cols[0]
        try:
            df[dcol] = pd.to_datetime(df[dcol], errors="coerce")
            valid_dates = df[dcol].dropna()
            if not valid_dates.empty:
                min_d = valid_dates.min().date()
                max_d = valid_dates.max().date()
                st.markdown("##### 分析日期區間")
                dr_col1, dr_col2 = st.columns(2)
                start_d = dr_col1.date_input("起始日期", value=min_d, min_value=min_d, max_value=max_d)
                end_d   = dr_col2.date_input("結束日期", value=max_d, min_value=min_d, max_value=max_d)
                df = df[(df[dcol].dt.date >= start_d) & (df[dcol].dt.date <= end_d)]
                st.caption(f"目前顯示 {len(df)} 筆 / 共 {len(df_full)} 筆")
        except Exception:
            pass

    # 組合 source_name = 日期區間（用於 PPT 封面）
    if start_d and end_d:
        ppt_source = f"{start_d.strftime('%Y/%m/%d')}～{end_d.strftime('%Y/%m/%d')}"
    else:
        ppt_source = st.session_state.get("source_name", "unknown")

    stats = df["問題類型"].value_counts().rename_axis("問題類型").reset_index(name="件數")
    stats["百分比"] = (stats["件數"] / max(stats["件數"].sum(), 1) * 100).round(0).astype(int)
    stats["歸屬部門"] = stats["問題類型"].map(DEPT_MAP).fillna("")

    # Build totals row
    total_count = int(stats["件數"].sum())
    dept_totals = stats.groupby("歸屬部門")["件數"].sum()
    dept_summary = "  ".join([f"{d}:{int(n)}件" for d, n in dept_totals.items() if d])
    totals_row = pd.DataFrame([{
        "問題類型": "[ 合計 ]",
        "件數": total_count,
        "百分比": 100,
        "歸屬部門": dept_summary,
    }])
    stats_with_total = pd.concat([stats, totals_row], ignore_index=True)

    st.markdown("#### 類型件數與部門 (可直接編輯，圖表即時同步)")
    edited_stats = st.data_editor(
        stats_with_total,
        use_container_width=True,
        hide_index=True,
        column_config={
            "歸屬部門": st.column_config.SelectboxColumn(options=DEPT_OPTIONS + [dept_summary]),
            "百分比": st.column_config.NumberColumn(format="%d %%")
        },
        key="stats_editor",
        num_rows="fixed",
    )
    # Use main stats (drop totals row) for charts
    chart_stats = edited_stats[edited_stats["問題類型"] != "[ 合計 ]"]
    render_charts_from_stats(chart_stats, df, key_prefix="sec2")

    st.markdown("#### AI 問題重點分析")
    # 摘要依「資料指紋」快取：同一份資料只呼叫一次 API，
    # 避免每次互動（Streamlit 會整頁重跑）都重新計費。
    ai_text = get_ai_summary_cached(df, st.session_state.get("source_name", ""))
    st.text_area("分析摘要預覽", ai_text, height=140)
    sum_c1, sum_c2 = st.columns([1.4, 6])
    if sum_c1.button("重新產生摘要", key="sec2_regen_summary"):
        st.session_state.pop(_ai_summary_cache_key(df, st.session_state.get("source_name", "")), None)
        st.rerun()
    if not (auto_config.get_anthropic_key() or auto_config.get_openai_key()):
        sum_c2.caption("尚未設定 ANTHROPIC_API_KEY／OPENAI_API_KEY，目前使用內建統計摘要。")

    # ── 預先產生所有下載檔案（避免 Streamlit on_click 時檔案還未產生）──
    chart_colors = st.session_state.get("chart_colors_sec2", {})

    # 用 session_state 快取，避免每次重繪都重新產生大檔
    cache_key = f"chart_pack_{ppt_source}"
    if cache_key not in st.session_state:
        with st.spinner("正在產生圖表與簡報..."):
            try:
                st.session_state[cache_key] = build_chart_pack(
                    df,
                    color_bar=chart_colors.get("bar"),
                    color_pie=chart_colors.get("pie"),
                    color_hbar=chart_colors.get("hbar"),
                )
            except Exception as e:
                st.error(f"圖表產生失敗：{e}")
                st.session_state[cache_key] = {}

    chart_pack = st.session_state[cache_key]

    ppt_cache_key = f"ppt_bytes_{ppt_source}"
    if ppt_cache_key not in st.session_state:
        with st.spinner("正在產生 PPT 簡報..."):
            try:
                st.session_state[ppt_cache_key] = build_ppt_bytes(
                    chart_stats, ai_text, ppt_source, chart_pack=chart_pack,
                )
            except Exception as e:
                st.error(f"PPT 產生失敗：{e}")
                st.session_state[ppt_cache_key] = b""

    ppt_bytes = st.session_state[ppt_cache_key]

    # ── 產生 ZIP ──
    zip_cache_key = f"zip_bytes_{ppt_source}"
    if zip_cache_key not in st.session_state:
        try:
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for fn, b in chart_pack.items():
                    zi = zipfile.ZipInfo(fn)
                    zi.flag_bits |= 0x800
                    zi.compress_type = zipfile.ZIP_DEFLATED
                    zf.writestr(zi, b)
            st.session_state[zip_cache_key] = zip_buf.getvalue()
        except Exception as e:
            st.error(f"ZIP 產生失敗：{e}")
            st.session_state[zip_cache_key] = b""

    zip_bytes = st.session_state[zip_cache_key]

    # ── 下載按鈕（檔案已預先備好）──
    dl_col1, dl_col2, dl_col3 = st.columns(3)
    dl_col1.download_button(
        "⬇️ 下載 AI 分析文字檔",
        data=ai_text.encode("utf-8"),
        file_name=f"{datetime.now().strftime('%Y%m%d')}_AI重點分析.txt",
        mime="text/plain",
        use_container_width=True,
    )
    dl_col2.download_button(
        "⬇️ 下載圖表圖檔（ZIP）",
        data=zip_bytes,
        file_name=f"{datetime.now().strftime('%Y%m%d')}_圖表圖檔.zip",
        mime="application/zip",
        use_container_width=True,
        disabled=not zip_bytes,
    )
    dl_col3.download_button(
        "⬇️ 一鍵下載分析簡報 PPT",
        data=ppt_bytes,
        file_name=f"{datetime.now().strftime('%Y%m%d')}_分析簡報.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        disabled=not ppt_bytes,
    )



def section_3():
    page_header("功能三：歷史分析紀錄", "歷史分析紀錄管理（最新置頂），可預覽與下載。")

    # ── Google Sheets 連線狀態 ──
    import os
    has_creds = bool(os.environ.get("GOOGLE_CREDENTIALS_JSON", ""))
    has_sid   = bool(os.environ.get("HISTORY_SHEET_ID", ""))
    ws_test   = _history_sheet()
    if ws_test is not None:
        st.success("☁️ Google Sheets 已連線，歷史紀錄永久保存")
    elif has_creds and has_sid:
        ws_test2 = _history_sheet(log_error=True)
        err_detail = st.session_state.get("_gsheet_error", "")
        st.warning(f"⚠️ 環境變數已設定但連線失敗\n{err_detail}")
        st.info("💡 請到 Google Cloud Console 確認已啟用 **Google Sheets API** 與 **Google Drive API**：\nhttps://console.cloud.google.com/apis/library")
    else:
        st.info("ℹ️ 未連線 Google Sheets，歷史紀錄僅限本次瀏覽")

    history = load_history()
    if not history:
        st.info("尚無歷史紀錄。")
        return

    # De-duplicate
    seen_names: dict = {}
    deduped = []
    for item in history:
        sn = item.get("source_name", "")
        if sn not in seen_names:
            seen_names[sn] = item
            deduped.append(item)
    history = deduped

    # ── 日期區間篩選器 ─────────────────────────────────────────
    st.markdown("---")
    st.markdown("##### 📅 日期區間篩選（有篩選時才顯示紀錄）")
    f_col1, f_col2, f_col3 = st.columns([2, 2, 1])

    # 取得所有紀錄的日期範圍
    all_dates = []
    for item in history:
        try:
            all_dates.append(datetime.fromisoformat(item["created_at"]).date())
        except Exception:
            pass

    if all_dates:
        min_date = min(all_dates)
        max_date = max(all_dates)
    else:
        min_date = max_date = datetime.now().date()

    start_filter = f_col1.date_input("開始日期", value=None, min_value=min_date, max_value=max_date,
                                      key="s3_start", format="YYYY/MM/DD")
    end_filter   = f_col2.date_input("結束日期", value=None, min_value=min_date, max_value=max_date,
                                      key="s3_end", format="YYYY/MM/DD")
    do_filter = f_col3.button("🔍 篩選", key="s3_filter", use_container_width=True)

    # 是否已啟動篩選
    filter_active = start_filter is not None or end_filter is not None

    if not filter_active:
        st.caption("請選擇日期區間後按「篩選」按鈕，即可顯示該區間的歷史紀錄。")
        return

    # 依日期區間過濾
    filtered = []
    for item in history:
        try:
            item_date = datetime.fromisoformat(item["created_at"]).date()
            if start_filter and item_date < start_filter:
                continue
            if end_filter and item_date > end_filter:
                continue
            filtered.append(item)
        except Exception:
            filtered.append(item)

    if not filtered:
        st.info(f"所選日期區間（{start_filter} ～ {end_filter}）無歷史紀錄。")
        return

    st.caption(f"共找到 {len(filtered)} 筆紀錄")
    history = filtered

    for item in history:
        out_path = Path(item.get("output_path", ""))
        cache = st.session_state.get("_history_cache", {})
        item_id = item["id"]

        # 取得 excel bytes：磁碟 → session_state 快取（已由 load_history 從 Sheets 填入）
        dl_bytes = None
        df_hist  = None
        if out_path.exists():
            try:
                dl_bytes = out_path.read_bytes()
                df_hist  = pd.read_excel(io.BytesIO(dl_bytes))
            except Exception:
                dl_bytes = None
        if dl_bytes is None and item_id in cache:
            try:
                dl_bytes = cache[item_id]["excel_bytes"]
                df_hist  = pd.read_excel(io.BytesIO(dl_bytes))
            except Exception:
                dl_bytes = None

        if dl_bytes is None:
            continue   # 真的找不到，跳過
        
        sname = item.get('source_name', '')
        if len(sname) > 28:
            sname = sname[:14] + "..." + sname[-10:]
        label = f"{item['created_at'][:16]}  {sname}  ({item['rows']} 筆)"
        with st.expander(label):
            tab_data, tab_chart, tab_ai = st.tabs(["資料預覽", "圖表分析", "AI 重點摘要"])
            
            with tab_data:
                st.dataframe(df_hist.head(30), use_container_width=True, hide_index=True)
                col1, col2, col3 = st.columns([1, 1, 1])
                col1.download_button(
                    "下載該分析檔",
                    data=dl_bytes,
                    file_name=item["output_name"],
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"download_{item['id']}",
                )
                if col2.button("[編輯]", key=f"edit_{item['id']}"):
                    st.session_state["analysis_df"] = df_hist.copy()
                    st.session_state["source_name"] = item["source_name"]
                    st.session_state["_editing_history_id"] = item["id"]
                    st.session_state["menu"] = "上傳檔案區（分析區）"
                    st.rerun()
                if col3.button("[刪除]", key=f"del_{item['id']}"):
                    delete_history(item["id"])
                    st.rerun()
            
            with tab_chart:
                if not df_hist.empty:
                    render_charts(df_hist, key_prefix=f"hist_{item['id']}")
                    cdl1, cdl2 = st.columns(2)
                    hist_stats = df_hist["問題類型"].value_counts().rename_axis("問題類型").reset_index(name="件數")
                    hist_stats["百分比"] = (hist_stats["件數"] / max(hist_stats["件數"].sum(), 1) * 100).round(0).astype(int)
                    hist_stats["歸屬部門"] = hist_stats["問題類型"].map(DEPT_MAP).fillna("")
                    hist_ai = generate_ai_summary(df_hist)
                    hist_chart_pack = build_chart_pack(df_hist)

                    hist_ppt = build_ppt_bytes(
                        hist_stats,
                        hist_ai,
                        item.get("source_name", "history"),
                        chart_pack=hist_chart_pack,
                    )
                    cdl1.download_button(
                        "一鍵下載PPT",
                        data=hist_ppt,
                        file_name=f"{datetime.now().strftime('%Y%m%d')}_{safe_filename(item.get('source_name','history'))}_圖表分析.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"hist_ppt_{item['id']}",
                    )
                    hist_zip = io.BytesIO()
                    with zipfile.ZipFile(hist_zip, "w", zipfile.ZIP_DEFLATED) as zf:
                        for fn, b in hist_chart_pack.items():
                            zi = zipfile.ZipInfo(fn)
                            zi.flag_bits |= 0x800  # UTF-8 filename flag，避免中文亂碼
                            zi.compress_type = zipfile.ZIP_DEFLATED
                            zf.writestr(zi, b)
                    cdl2.download_button(
                        "下載圖檔（ZIP）",
                        data=hist_zip.getvalue(),
                        file_name=f"{datetime.now().strftime('%Y%m%d')}_{safe_filename(item.get('source_name','history'))}_圖表.zip",
                        mime="application/zip",
                        key=f"hist_img_{item['id']}",
                    )
                else:
                    st.info("無資料可繪圖")
                    
            with tab_ai:
                st.info("點擊下方按鈕即時生成本檔案的 AI 重點摘要")
                if st.button("[產生 AI 摘要]", key=f"ai_btn_{item['id']}"):
                    with st.spinner("AI 分析中..."):
                        ai_result = generate_ai_summary_llm(df_hist)
                        st.markdown(ai_result)



def section_4():
    """功能四：週/月/季/年度趨勢分析儀表板"""

    # ── ECOCO 品牌 CSS（對齊 HTML 範本風格）──────────────────────
    st.markdown("""<style>
    .s4-section{border-left:6px solid #FF5000;padding-left:14px;
                color:#060E9F;margin:24px 0 14px}
    [data-testid="stAppViewContainer"] .s4-section,
    .s4-section{font-size:26px !important;font-weight:700 !important;line-height:1.3}
    .s4-card{background:#fff;border-radius:12px;padding:20px 24px;
             box-shadow:0 4px 10px rgba(0,0,0,.06);margin-bottom:16px}
    .s4-kpi-grid{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin-bottom:18px}
    .s4-kpi{background:#fff;border-radius:10px;padding:18px 16px;text-align:center;
            border-top:4px solid #FFCE00;box-shadow:0 2px 6px rgba(0,0,0,.05)}
    .s4-kpi-val{font-size:30px;font-weight:700;color:#FF5000}
    .s4-kpi-lbl{font-size:12px;color:#666;margin-top:4px}
    .s4-kpi-delta{font-size:11px;margin-top:3px}
    .delta-up{color:#c03000} .delta-dn{color:#0a6e44} .delta-flat{color:#888}
    .s4-rank-table{width:100%;border-collapse:collapse}
    .s4-rank-table th{background:#060E9F;color:#fff;padding:10px 14px;font-size:14px;text-align:center;font-weight:600}
    .s4-rank-table td{padding:9px 14px;text-align:center;border-bottom:1px solid #eee;font-size:12px}
    .s4-rank-val{color:#FF5000;font-weight:700;font-size:14px}
    .filter-chip-row{display:flex;gap:8px;flex-wrap:wrap;margin-bottom:14px}
    .filter-chip{padding:4px 14px;border-radius:20px;font-size:12px;font-weight:600;cursor:pointer;
                 border:1.5px solid #060E9F;background:#fff;color:#060E9F}
    .filter-chip.active{background:#060E9F;color:#fff}
    </style>""", unsafe_allow_html=True)

    # ── 頁首 ──────────────────────────────────────────────────────
    page_header("📈 ECOCO 客訴趨勢分析儀表板",
                "城市・站點・部門・問題類型・機台比例 | 自訂日期區間 + 維度篩選")

    # ── 資料來源 ──────────────────────────────────────────────────
    src_tab1, src_tab2 = st.tabs(["📂 歷史紀錄資料", "🔗 填入 Google Sheets 網址"])
    all_dfs: list[pd.DataFrame] = []

    with src_tab1:
        ws = _history_sheet()
        if ws:
            import base64 as _b64
            try:
                for grow in ws.get_all_values()[1:]:
                    if not grow or not grow[0]: continue
                    data_ref = grow[4] if len(grow) > 4 else ""
                    if data_ref:
                        try:
                            if data_ref.startswith("sheet:"):
                                data_ws = ws.spreadsheet.worksheet(data_ref.split(":", 1)[1])
                                all_dfs.append(_worksheet_to_dataframe(data_ws))
                            else:
                                all_dfs.append(pd.read_excel(io.BytesIO(_b64.b64decode(data_ref))))
                        except Exception: pass
            except Exception: pass
        st.caption(f"已載入 {len(all_dfs)} 份歷史紀錄" if all_dfs else "尚無歷史資料")

    with src_tab2:
        gs_url = st.text_input("Google Sheets 網址", placeholder="https://docs.google.com/spreadsheets/d/xxxxx/edit", key="s4v3_gsurl")
        gs_sheet = st.text_input("工作表名稱（留空讀取第一張）", key="s4v3_gssheet", value="")
        if st.button("📥 讀取", key="s4v3_load_gs"):
            if not gs_url:
                st.error("請填入網址")
            else:
                try:
                    import re as _re
                    m = _re.search(r"/spreadsheets/d/([^/]+)", gs_url)
                    if not m:
                        st.error("無法解析試算表 ID")
                    else:
                        _client = _get_gsheet_client()
                        if not _client:
                            st.error("未連線 Google API")
                        else:
                            _ss = _client.open_by_key(m.group(1))
                            _ws = _ss.worksheet(gs_sheet) if gs_sheet else _ss.get_worksheet(0)
                            _rows = _ws.get_all_values()
                            if _rows:
                                _df = pd.DataFrame(_rows[1:], columns=_rows[0])
                                all_dfs.append(_df)
                                st.session_state["_s4v3_gs_df"] = _df
                                st.success(f"✅ 已讀取「{_ws.title}」，共 {len(_df)} 列")
                except Exception as e:
                    st.warning("讀取失敗，請確認網址、工作表名稱與試算表共用權限。")
        if st.session_state.get("_s4v3_gs_df") is not None:
            all_dfs.append(st.session_state["_s4v3_gs_df"])

    if not all_dfs:
        empty_state("沒有資料紀錄　—　請先在功能一完成分析儲存，或填入 Google Sheets 網址。")
        return

    # 合併前確保每份 df 欄位名稱唯一（避免重複欄位造成 InvalidIndexError）
    clean_dfs = []
    for _d in all_dfs:
        try:
            _d = _d.copy()
            # 若有重複欄位名稱，加後綴區分
            _seen = {}
            _new_cols = []
            for c in _d.columns:
                if c in _seen:
                    _seen[c] += 1
                    _new_cols.append(f"{c}_{_seen[c]}")
                else:
                    _seen[c] = 0
                    _new_cols.append(c)
            _d.columns = _new_cols
            clean_dfs.append(_d)
        except Exception:
            clean_dfs.append(_d)

    df_all = pd.concat(clean_dfs, ignore_index=True)
    # drop_duplicates 要求 index 唯一，先 reset_index
    try:
        df_all = df_all.loc[:, ~df_all.columns.duplicated()]  # 移除重複欄位
        df_all = df_all.drop_duplicates().reset_index(drop=True)
    except Exception:
        df_all = df_all.reset_index(drop=True)

    # 舊歷史紀錄可能是遮蔽功能上線前存的，顯示前補遮一次
    df_all = mask_phone_columns(df_all)

    # ── 欄位自動偵測 ──────────────────────────────────────────────
    date_col   = next((c for c in df_all.columns if "日期" in c or "date" in c.lower()), None)
    type_col   = next((c for c in df_all.columns if "問題類型" in c), None)
    detail_col = next((c for c in df_all.columns if "問題細項" in c), None)
    dept_col   = next((c for c in df_all.columns if "部門" in c or "歸屬" in c), None)
    city_col   = next((c for c in df_all.columns if "站點區域" in c or "城市" in c or "區域" in c), None)
    station_col= next((c for c in df_all.columns if c == "站點名稱"), None) or \
                 next((c for c in df_all.columns if "站點名稱" in c and "編號" not in c), None)
    machine_col= next((c for c in df_all.columns if "機台類型" in c or "機台" in c), None)

    if not date_col:
        empty_state("沒有資料紀錄　—　來源資料找不到日期欄位。")
        return

    df_all[date_col] = pd.to_datetime(df_all[date_col], errors="coerce")
    df_all = df_all.dropna(subset=[date_col])
    if df_all.empty:
        empty_state("沒有資料紀錄　—　來源資料沒有可辨識的日期。")
        return

    # ── 時間區間選擇（維度 + 自訂日期）────────────────────────────
    st.markdown('<div class="s4-section">⚙️ 篩選條件</div>', unsafe_allow_html=True)
    filter_c1, filter_c2, filter_c3 = st.columns([2, 3, 2])

    dim_mode = filter_c1.radio("時間模式", ["維度選擇", "自訂日期區間"], horizontal=True, key="s4v3_dimmode")

    DIM_FREQ = {"週": "W", "月": "M", "季": "Q", "年度": "Y"}
    period_sel = period_prev = None
    df_cur = df_prev = pd.DataFrame()

    if dim_mode == "維度選擇":
        dim = filter_c2.selectbox("分析維度", ["週", "月", "季", "年度"], index=1, key="s4v3_dim")
        df_all["_period"] = df_all[date_col].dt.to_period(DIM_FREQ[dim]).astype(str)
        periods = sorted(df_all["_period"].unique(), reverse=True)
        if not periods:
            empty_state()
            return
        period_sel = filter_c3.selectbox(f"本期", periods, key="s4v3_period")
        p_idx = periods.index(period_sel)
        period_prev = periods[p_idx + 1] if p_idx + 1 < len(periods) else None
        df_cur  = df_all[df_all["_period"] == period_sel].copy()
        df_prev = df_all[df_all["_period"] == period_prev].copy() if period_prev else pd.DataFrame()
        period_label = period_sel
    else:
        min_d = df_all[date_col].min().date()
        max_d = df_all[date_col].max().date()
        d_col1, d_col2 = filter_c2.columns(2)
        start_d = d_col1.date_input("開始", value=min_d, min_value=min_d, max_value=max_d, key="s4v3_sd")
        end_d   = d_col2.date_input("結束", value=max_d, min_value=min_d, max_value=max_d, key="s4v3_ed")
        df_cur  = df_all[(df_all[date_col].dt.date >= start_d) & (df_all[date_col].dt.date <= end_d)].copy()
        period_label = f"{start_d} ～ {end_d}"
        enable_compare = filter_c3.checkbox("啟用對照期", value=False, key="s4v3_compare_on")
        if enable_compare:
            cmp_c1, cmp_c2 = filter_c3.columns(2)
            cmp_start = cmp_c1.date_input("對照開始", value=min_d, min_value=min_d, max_value=max_d, key="s4v3_cmp_sd")
            cmp_end = cmp_c2.date_input("對照結束", value=min_d, min_value=min_d, max_value=max_d, key="s4v3_cmp_ed")
            df_prev = df_all[(df_all[date_col].dt.date >= cmp_start) & (df_all[date_col].dt.date <= cmp_end)].copy()
            period_prev = f"{cmp_start} ～ {cmp_end}"
        else:
            period_prev = None
            df_prev = pd.DataFrame()

    # ── 多維篩選 chips（城市/部門/問題類型/機台）──────────────────
    st.markdown("**篩選維度：**")
    chip_cols = st.columns(4)

    city_filter   = chip_cols[0].multiselect("🏙️ 城市", sorted(df_cur[city_col].dropna().unique().tolist()) if city_col and city_col in df_cur.columns else [], key="s4v3_city")
    dept_filter   = chip_cols[1].multiselect("🏢 部門", sorted(df_cur[dept_col].dropna().unique().tolist()) if dept_col and dept_col in df_cur.columns else [], key="s4v3_dept")
    type_filter   = chip_cols[2].multiselect("❓ 問題類型", sorted(df_cur[type_col].dropna().unique().tolist()) if type_col and type_col in df_cur.columns else [], key="s4v3_type")
    mach_filter   = chip_cols[3].multiselect("🔧 機台類型", sorted(df_cur[machine_col].dropna().unique().tolist()) if machine_col and machine_col in df_cur.columns else [], key="s4v3_mach")

    df_filt = df_cur.copy()
    if city_filter  and city_col:   df_filt = df_filt[df_filt[city_col].isin(city_filter)]
    if dept_filter  and dept_col:   df_filt = df_filt[df_filt[dept_col].isin(dept_filter)]
    if type_filter  and type_col:   df_filt = df_filt[df_filt[type_col].isin(type_filter)]
    if mach_filter  and machine_col: df_filt = df_filt[df_filt[machine_col].isin(mach_filter)]

    n_cur  = len(df_filt)
    n_prev = len(df_prev)

    def pct_change(cur, prev):
        if prev == 0: return None
        return (cur - prev) / prev * 100

    # ── 機台類型：將「方舟」歸類為「收瓶機」────────────────────────
    if machine_col and machine_col in df_filt.columns:
        def _normalize_machine(val):
            v = str(val).strip()
            if "方舟" in v or "收瓶" in v: return "收瓶機"
            if "電池" in v: return "電池機"
            return v
        df_filt = df_filt.copy()
        df_filt[machine_col] = df_filt[machine_col].apply(_normalize_machine)
    if machine_col and machine_col in df_all.columns:
        df_all = df_all.copy()
        df_all[machine_col] = df_all[machine_col].apply(
            lambda v: "收瓶機" if ("方舟" in str(v) or "收瓶" in str(v)) else ("電池機" if "電池" in str(v) else str(v))
        )

    # 篩選後沒有資料就到此為止，後面的 KPI／排行／圖表都沒有東西可畫，
    # 硬跑下去只會噴出一堆紅色錯誤。
    if df_filt.empty:
        st.markdown(f'<div class="s4-section">📊 本期即時統計（{period_label}）</div>',
                    unsafe_allow_html=True)
        empty_state("沒有資料紀錄")
        st.caption("請放寬日期區間或取消部分篩選條件。")
        return

    # ── KPI 卡片（用 st.metric 避免 HTML escape 問題）────────────────
    st.markdown(f'<div class="s4-section">📊 本期即時統計（{period_label}）</div>', unsafe_allow_html=True)
    st.caption(f"📅 資料區間：{period_label}　篩選後共 **{n_cur}** 筆")

    kpi_items = [("🗂️ 總進件數", n_cur, n_prev)]
    if type_col and type_col in df_filt.columns:
        for t, tc in df_filt[type_col].value_counts().items():
            prev_tc = int(df_prev[type_col].eq(t).sum()) if not df_prev.empty and type_col in df_prev.columns else 0
            kpi_items.append((str(t)[:12], int(tc), prev_tc))
            if len(kpi_items) >= 4: break

    kpi_cols = st.columns(len(kpi_items[:4]))
    for col_i, (lbl, cur, prev) in enumerate(kpi_items[:4]):
        p = pct_change(cur, prev)
        delta_str = None
        if p is not None:
            sym = "+" if p >= 0 else ""
            delta_str = f"{sym}{p:.1f}% vs 上期"
        kpi_cols[col_i].metric(label=lbl, value=cur, delta=delta_str)

    # ── 排行統計（區域/站點/問題細項）───────────────────────────────
    st.markdown(f'<div class="s4-section">🏆 案件排行統計 Top 5 ── {period_label}</div>', unsafe_allow_html=True)

    rank_cols = st.columns(3)
    MEDAL = ["🥇","🥈","🥉","4️⃣","5️⃣"]

    def rank_table_html(series, header1, header2):
        rows = ""
        for idx, (k, v) in enumerate(series.head(5).items()):
            m = MEDAL[idx] if idx < len(MEDAL) else str(idx+1)
            rows += (f'<tr><td style="text-align:left;font-size:12px">{m} {str(k)[:24]}</td>'
                     f'<td class="s4-rank-val">{int(v)}</td></tr>')
        return f'''<table class="s4-rank-table">
          <thead><tr><th>{header1}</th><th>{header2}</th></tr></thead>
          <tbody>{rows}</tbody>
        </table>'''

    with rank_cols[0]:
        st.markdown('<div style="font-weight:700;color:#060E9F;margin-bottom:8px">📍 區域排行</div>', unsafe_allow_html=True)
        if city_col and city_col in df_filt.columns and not df_filt[city_col].dropna().empty:
            st.markdown('<div class="s4-card">' + rank_table_html(df_filt[city_col].value_counts(), "城市/區域", "件數") + '</div>', unsafe_allow_html=True)
        else:
            empty_state()

    with rank_cols[1]:
        st.markdown('<div style="font-weight:700;color:#060E9F;margin-bottom:8px">🏬 站點排行</div>', unsafe_allow_html=True)
        if station_col and station_col in df_filt.columns and not df_filt[station_col].dropna().empty:
            st.markdown('<div class="s4-card">' + rank_table_html(df_filt[station_col].value_counts(), "站點名稱", "件數") + '</div>', unsafe_allow_html=True)
        else:
            empty_state()

    with rank_cols[2]:
        st.markdown('<div style="font-weight:700;color:#060E9F;margin-bottom:8px">🔍 問題細項排行</div>', unsafe_allow_html=True)
        if detail_col and detail_col in df_filt.columns and not df_filt[detail_col].dropna().empty:
            st.markdown('<div class="s4-card">' + rank_table_html(df_filt[detail_col].value_counts(), "問題細項", "件數") + '</div>', unsafe_allow_html=True)
        else:
            empty_state()

    # ── 圖表：問題類型 + 機台佔比（對齊 HTML 範本）─────────────────
    st.markdown(f'<div class="s4-section">📉 數據可視化分析 ── {period_label}</div>', unsafe_allow_html=True)
    chart_col1, chart_col2 = st.columns(2)

    with chart_col1:
        if type_col and type_col in df_filt.columns:
            _tc = df_filt[type_col].value_counts()
            _total = _tc.sum()
            COLORS_PIE = ["#060E9F","#FF5000","#FFCE00","#8EB9C9","#0076A9","#FAE0B8"]
            # 圖例文字直接寫進 names：圓餅圖只有一個 trace，
            # 用 for_each_trace 改 name 是改不到圖例的（圖例讀的是 label）。
            fig_pie = px.pie(
                values=_tc.values, names=pie_legend_labels(_tc),
                title=f"{period_label} 客訴類別分佈",
                hole=0.38,
                color_discrete_sequence=COLORS_PIE,
            )
            fig_pie.update_traces(
                texttemplate="%{percent:.0%}",   # 只在扇形內顯示 %
                textposition="inside",
                textfont=dict(size=13, color="white"),
                hovertemplate="<b>%{label}</b><br>%{value}件 / %{percent:.1%}<extra></extra>",
                showlegend=True,
            )
            fig_pie.update_layout(height=380, showlegend=True,
                                  legend=PIE_LEGEND, margin=PIE_MARGIN,
                                  title_font_size=14, title_x=0.0)
            st.plotly_chart(fig_pie, use_container_width=True)

    with chart_col2:
        if machine_col and machine_col in df_filt.columns and not df_filt[machine_col].dropna().empty:
            _mc = df_filt[machine_col].value_counts()
            _mc_total = _mc.sum()
            COLORS_MAC = ["#FF5000","#060E9F","#8EB9C9","#FFCE00"]
            fig_mac = px.pie(
                values=_mc.values, names=pie_legend_labels(_mc),
                title=f"{period_label} 機台客訴佔比",
                color_discrete_sequence=COLORS_MAC,
            )
            fig_mac.update_traces(
                texttemplate="%{percent:.0%}",
                textposition="inside",
                textfont=dict(size=14, color="white"),
                hovertemplate="<b>%{label}</b><br>%{value}件 / %{percent:.1%}<extra></extra>",
            )
            fig_mac.update_layout(height=380, showlegend=True,
                                  legend=PIE_LEGEND, margin=PIE_MARGIN,
                                  title_font_size=14)
            st.plotly_chart(fig_mac, use_container_width=True)
        elif detail_col and detail_col in df_filt.columns:
            _dc = df_filt[detail_col].value_counts().head(8)
            fig_det = px.bar(
                x=list(_dc.values)[::-1], y=list(_dc.index)[::-1],
                orientation="h", title=f"{period_label} TOP 8 問題細項",
                color_discrete_sequence=["#060E9F"],
            )
            fig_det.update_layout(height=420, xaxis=dict(tickformat="d", nticks=6),
                                   margin=dict(t=45,b=0,l=0,r=0))
            st.plotly_chart(fig_det, use_container_width=True)

    # ── 趨勢折線圖 ────────────────────────────────────────────────
    st.markdown(f'<div class="s4-section">📈 客訴趨勢分析 ── {period_label}</div>', unsafe_allow_html=True)
    if dim_mode == "維度選擇" and len(df_all["_period"].unique()) >= 2:
        _trend = df_all.groupby("_period").size().reset_index(name="件數").sort_values("_period")
        fig_line = px.line(
            _trend, x="_period", y="件數",
            title=f"歷史件數趨勢",
            markers=True,
            color_discrete_sequence=["#FF5000"],
        )
        fig_line.update_traces(fill="tozeroy", fillcolor="rgba(255,80,0,0.1)")
        if period_sel and period_sel in _trend["_period"].values:
            _sel_i = _trend.index[_trend["_period"] == period_sel].tolist()
            if _sel_i:
                fig_line.add_vline(x=_sel_i[0], line_dash="dash", line_color="#060E9F",
                                   annotation_text="本期", annotation_font_color="#060E9F")
        fig_line.update_layout(
            height=320, xaxis_title="期間",
            yaxis=dict(tickformat="d", nticks=6),
            paper_bgcolor="white", plot_bgcolor="rgba(250,224,184,0.15)",
            margin=dict(t=45,b=0),
        )
        st.plotly_chart(fig_line, use_container_width=True)
    else:
        # 自訂日期：每日件數
        _daily = df_filt.groupby(df_filt[date_col].dt.date).size().reset_index(name="件數")
        _daily.columns = ["日期", "件數"]
        if not _daily.empty:
            fig_daily = px.bar(
                _daily, x="日期", y="件數",
                title="期間內每日件數",
                color_discrete_sequence=["#060E9F"],
            )
            fig_daily.update_layout(height=300, yaxis=dict(tickformat="d", nticks=6), margin=dict(t=45,b=0))
            st.plotly_chart(fig_daily, use_container_width=True)

    # ── 城市展開排行（可折疊）────────────────────────────────────
    if city_col and city_col in df_filt.columns and not df_filt.empty:
        st.markdown(f'<div class="s4-section">🏙️ 區域排行榜 ── {period_label}</div>', unsafe_allow_html=True)
        city_rank = df_filt[city_col].value_counts()
        MEDAL_LIST = ["🥇","🥈","🥉","4️⃣","5️⃣","6️⃣","7️⃣","8️⃣","9️⃣","🔟"]
        for ri, (city, cnt) in enumerate(city_rank.items()):
            prev_cnt = int(df_prev[city_col].eq(city).sum()) if not df_prev.empty and city_col in df_prev.columns else 0
            p = pct_change(int(cnt), prev_cnt)
            delta_s = (f"　{'▲' if p>0 else '▼'}{abs(p):.1f}%") if p is not None else ""
            medal = MEDAL_LIST[ri] if ri < len(MEDAL_LIST) else f"#{ri+1}"
            with st.expander(f"{medal} **{city}**　{int(cnt)} 件{delta_s}", expanded=(ri==0)):
                df_city = df_filt[df_filt[city_col] == city]
                ec1, ec2 = st.columns(2)
                with ec1:
                    st.markdown("**📍 站點排行**")
                    if station_col and station_col in df_city.columns:
                        _sr = df_city[station_col].value_counts().head(8)
                        _sr_prev = df_prev[df_prev[city_col]==city][station_col].value_counts() if not df_prev.empty and city_col in df_prev.columns and station_col in df_prev.columns else pd.Series(dtype=int)
                        _html = ""
                        for si, (sn, sv) in enumerate(_sr.items()):
                            _sp = int(_sr_prev.get(sn, 0))
                            _sd = (f"　{'▲' if int(sv)-_sp>0 else '▼'}{abs(pct_change(int(sv),_sp)):.0f}%") if _sp and pct_change(int(sv),_sp) is not None else ""
                            _sm = MEDAL_LIST[si] if si < len(MEDAL_LIST) else f"#{si+1}"
                            _html += f'<div style="padding:5px 0;border-bottom:.5px solid #eee;font-size:13px;display:flex;justify-content:space-between"><span>{_sm} {str(sn)[:18]}</span><b style="color:#FF5000">{int(sv)}{_sd}</b></div>'
                        st.markdown(_html, unsafe_allow_html=True)
                with ec2:
                    st.markdown("**🔍 問題細項排行**")
                    if detail_col and detail_col in df_city.columns:
                        _dr = df_city[detail_col].value_counts().head(8)
                        _max = int(_dr.max()) if not _dr.empty else 1
                        _dhtml = ""
                        for di, (dn, dv) in enumerate(_dr.items()):
                            _dm = MEDAL_LIST[di] if di < len(MEDAL_LIST) else f"#{di+1}"
                            _bp = int(dv)/_max*100
                            _dhtml += f'''<div style="padding:5px 0;border-bottom:.5px solid #eee">
                              <div style="display:flex;justify-content:space-between;font-size:12px;margin-bottom:3px">
                                <span>{_dm} {str(dn)[:14]}</span><b style="color:#060E9F">{int(dv)}</b></div>
                              <div style="background:#f0f0f0;border-radius:3px;height:5px">
                                <div style="background:#060E9F;width:{_bp:.0f}%;height:100%;border-radius:3px"></div></div>
                            </div>'''
                        st.markdown(_dhtml, unsafe_allow_html=True)

    # ── 部門分析 ─────────────────────────────────────────────────
    if dept_col and dept_col in df_filt.columns:
        st.markdown(f'<div class="s4-section">🏢 各部門件數分析 ── {period_label}</div>', unsafe_allow_html=True)
        dept_rank = df_filt[dept_col].replace("","未分配").value_counts()
        DEPT_COLOR = {"營運部":"#FF5000","行銷部":"#FFCE00","資訊部":"#060E9F"}
        fig_dept = px.bar(
            dept_rank.reset_index(), x=dept_col, y="count",
            title="各部門件數",
            color=dept_col,
            color_discrete_map=DEPT_COLOR,
        )
        fig_dept.update_layout(height=300, yaxis=dict(tickformat="d", nticks=6),
                                showlegend=False, margin=dict(t=45,b=0))
        st.plotly_chart(fig_dept, use_container_width=True)

    # ── 完整頁面 PDF 下載 ────────────────────────────────────────
    st.markdown("---")
    st.markdown(f'<div class="s4-section">⬇️ 下載完整分析報告</div>', unsafe_allow_html=True)

    if st.button("📄 產生完整分析 PDF", key="s4_full_pdf", use_container_width=False):
        with st.spinner("正在產生多頁 PDF 報告..."):
            try:
                from fpdf import FPDF
                from fpdf.enums import XPos, YPos
                import os, glob, matplotlib.pyplot as _mplt
                from matplotlib.ticker import MaxNLocator as _MNL

                _FONT_CANDS = [
                    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/opentype/noto/NotoSansCJK-Medium.ttc",
                    "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
                    "/usr/share/fonts/truetype/arphic/uming.ttc",
                    "/tmp/NotoSansCJK.ttc",
                ]
                _FONT_CANDS += glob.glob("/usr/share/fonts/**/NotoSansCJK*.ttc", recursive=True)
                _font_path = _ensure_cjk_font()  # 使用已快取的字型路徑
                if not _font_path:
                    _font_path = next((p for p in _FONT_CANDS if os.path.exists(p)), None)

                _setup_cjk_font()

                class EcocoPDF(FPDF):
                    def __init__(self, font_path, font_name):
                        super().__init__(orientation="P", format="A4")
                        self.fp = font_path
                        self.fn = font_name
                        self.set_auto_page_break(auto=True, margin=15)
                        if font_path:
                            self.add_font(font_name, style="", fname=font_path)
                        self.set_margins(15, 15, 15)

                    def header(self):
                        # 藍色頁首條
                        self.set_fill_color(6, 14, 159)
                        self.rect(0, 0, 210, 14, style="F")
                        self.set_font(self.fn, size=9)
                        self.set_text_color(255, 255, 255)
                        self.set_xy(5, 3)
                        self.cell(0, 8, self._s(f"ECOCO 客訴趨勢分析報告　{period_label}"))
                        self.set_draw_color(255, 80, 0)
                        self.set_line_width(1.5)
                        self.line(0, 14, 210, 14)
                        self.set_line_width(0.2)
                        self.set_text_color(30, 30, 30)
                        self.ln(6)

                    def footer(self):
                        self.set_y(-12)
                        self.set_font(self.fn, size=8)
                        self.set_text_color(150, 150, 150)
                        self.cell(0, 8, self._s(f"第 {self.page_no()} 頁　產出日期：{datetime.now().strftime('%Y/%m/%d')}"), align="C")

                    def _s(self, s):
                        return s if self.fn != "Helvetica" else s.encode("ascii", "replace").decode()

                    def section_title(self, title):
                        self.ln(3)
                        self.set_fill_color(255, 80, 0)
                        self.rect(self.get_x(), self.get_y(), 6, 9, style="F")
                        self.set_font(self.fn, size=13)
                        self.set_text_color(6, 14, 159)
                        self.set_x(self.get_x() + 9)
                        self.cell(0, 9, self._s(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                        self.ln(2)

                    def full_table(self, headers, rows, col_widths):
                        """全幅表格，支援多行自動換行"""
                        # 表頭
                        self.set_fill_color(6, 14, 159)
                        self.set_text_color(255, 255, 255)
                        self.set_font(self.fn, size=9)
                        x0 = self.get_x()
                        for h, w in zip(headers, col_widths):
                            self.cell(w, 8, self._s(str(h)), border=1, fill=True, align="C",
                                      new_x=XPos.RIGHT, new_y=YPos.TOP)
                        self.ln(8)
                        # 資料列
                        self.set_text_color(30, 30, 30)
                        self.set_font(self.fn, size=9)
                        for i, row in enumerate(rows):
                            bg = (235, 244, 250) if i % 2 == 0 else (255, 255, 255)
                            self.set_fill_color(*bg)
                            self.set_x(x0)
                            for val, w in zip(row, col_widths):
                                self.cell(w, 7, self._s(str(val)[:40]), border=1, fill=True,
                                          new_x=XPos.RIGHT, new_y=YPos.TOP)
                            self.ln(7)
                        self.ln(3)

                    def embed_image(self, fig, w=180, h=110):
                        """嵌入 matplotlib 圖表"""
                        _b = io.BytesIO()
                        fig.savefig(_b, format="png", dpi=180, bbox_inches="tight",
                                    facecolor="white")
                        _mplt.close(fig)
                        _b.seek(0)
                        x = (210 - w) / 2  # 置中
                        self.image(_b, x=x, y=self.get_y(), w=w, h=h)
                        self.set_y(self.get_y() + h + 4)

                F = "CJK" if _font_path else "Helvetica"

                pdf = EcocoPDF(_font_path, F)

                # ════════════════════════════════════════════
                # Page 1：封面 + KPI 摘要
                # ════════════════════════════════════════════
                pdf.add_page()
                # 大標題框
                pdf.set_fill_color(6, 14, 159)
                pdf.rect(15, 20, 180, 38, style="F")
                pdf.set_font(F, size=20)
                pdf.set_text_color(255, 255, 255)
                pdf.set_xy(15, 26)
                pdf.cell(180, 12, pdf._s("ECOCO 客訴趨勢分析報告"), align="C",
                         new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                pdf.set_font(F, size=11)
                pdf.set_xy(15, 42)
                pdf.cell(180, 8, pdf._s(f"資料區間：{period_label}"), align="C",
                         new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                pdf.set_draw_color(255, 80, 0)
                pdf.set_line_width(2)
                pdf.line(15, 58, 195, 58)
                pdf.set_line_width(0.2)
                pdf.set_y(68)

                # KPI 卡片（橫排）
                pdf.section_title("本期即時統計摘要")
                kpi_data = [("🗂️ 總進件數", n_cur)]
                if type_col and type_col in df_filt.columns:
                    for t, tc in df_filt[type_col].value_counts().head(3).items():
                        kpi_data.append((str(t), int(tc)))
                card_w = 170 // len(kpi_data)
                card_x = 20
                for lbl, val in kpi_data:
                    pdf.set_fill_color(255, 206, 0)  # 黃色頂線
                    pdf.rect(card_x, pdf.get_y(), card_w - 4, 3, style="F")
                    pdf.set_fill_color(248, 249, 252)
                    pdf.rect(card_x, pdf.get_y() + 3, card_w - 4, 28, style="F")
                    pdf.set_font(F, size=22)
                    pdf.set_text_color(255, 80, 0)
                    pdf.set_xy(card_x, pdf.get_y() + 5)
                    pdf.cell(card_w - 4, 14, str(val), align="C",
                             new_x=XPos.RIGHT, new_y=YPos.TOP)
                    pdf.set_font(F, size=8)
                    pdf.set_text_color(80, 80, 80)
                    pdf.set_xy(card_x, pdf.get_y() + 20)
                    pdf.cell(card_w - 4, 8, pdf._s(str(lbl)[:10]), align="C",
                             new_x=XPos.RIGHT, new_y=YPos.TOP)
                    card_x += card_w
                pdf.set_y(pdf.get_y() + 36)

                # 城市 KPI
                if city_col and city_col in df_filt.columns:
                    pdf.ln(4)
                    pdf.set_font(F, size=9)
                    pdf.set_text_color(60, 60, 60)
                    city_top = df_filt[city_col].value_counts().head(3)
                    line = "　|　".join([f"{c}：{int(v)} 件" for c, v in city_top.items()])
                    pdf.set_x(20)
                    pdf.cell(0, 7, pdf._s(f"前三大城市：{line}"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

                # ════════════════════════════════════════════
                # Page 2：排行統計
                # ════════════════════════════════════════════
                pdf.add_page()
                pdf.section_title("案件排行統計")

                # 區域排行
                if city_col and city_col in df_filt.columns:
                    pdf.set_font(F, size=10); pdf.set_text_color(6,14,159)
                    pdf.cell(0, 7, pdf._s("📍 城市/區域排行"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                    _city_v = df_filt[city_col].value_counts()
                    rows_c = [[i+1, c, int(v), f"{int(v)/n_cur*100:.0f}%"]
                               for i, (c, v) in enumerate(_city_v.head(10).items())]
                    pdf.full_table(["排名","城市/區域","件數","佔比"], rows_c, [15,110,25,30])

                # 站點排行
                if station_col and station_col in df_filt.columns:
                    pdf.set_font(F, size=10); pdf.set_text_color(6,14,159)
                    pdf.cell(0, 7, pdf._s("🏬 站點排行"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                    _sta_v = df_filt[station_col].value_counts()
                    rows_s = [[i+1, str(s)[:30], int(v)]
                               for i, (s, v) in enumerate(_sta_v.head(10).items())]
                    pdf.full_table(["排名","站點名稱","件數"], rows_s, [15,140,25])

                # ════════════════════════════════════════════
                # Page 3：問題細項排行
                # ════════════════════════════════════════════
                pdf.add_page()
                pdf.section_title("問題細項排行")

                if detail_col and detail_col in df_filt.columns:
                    _det_v = df_filt[detail_col].value_counts()
                    rows_d = [[i+1, str(d)[:35], int(v), f"{int(v)/n_cur*100:.0f}%"]
                               for i, (d, v) in enumerate(_det_v.head(15).items())]
                    pdf.full_table(["排名","問題細項","件數","佔比"], rows_d, [15,120,20,25])

                # ════════════════════════════════════════════
                # Page 4：圖表（圓餅 + 機台）
                # ════════════════════════════════════════════
                pdf.add_page()
                pdf.section_title("數據可視化分析")

                if type_col and type_col in df_filt.columns:
                    _tc4 = df_filt[type_col].value_counts()
                    _total4 = _tc4.sum()
                    # 大圖（figsize 更寬），標籤改用圖例避免重疊
                    _f4, _a4 = _mplt.subplots(figsize=(9, 6))
                    _clrs4 = ["#060E9F","#FF5000","#FFCE00","#8EB9C9","#0076A9","#FAE0B8"]
                    _labels4 = [f"{k}（{int(v)}件）" for k, v in _tc4.items()]
                    wedges, texts, autotexts = _a4.pie(
                        list(_tc4.values),
                        labels=None,           # 不在扇形上顯示標籤，改用圖例
                        autopct=lambda p: f"{p:.0f}%" if p >= 5 else "",  # 小扇形不顯示%
                        colors=_clrs4[:len(_tc4)],
                        startangle=90,
                        pctdistance=0.75,
                        wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
                    )
                    for _at in autotexts:
                        _at.set_fontsize(11)
                        _at.set_fontweight("bold")
                    _a4.legend(
                        wedges, _labels4,
                        loc="center left",
                        bbox_to_anchor=(1.0, 0.5),
                        fontsize=9,
                        frameon=False,
                    )
                    _a4.set_title(f"{period_label}　客訴類別分佈", fontsize=13, pad=12)
                    _f4.tight_layout()
                    pdf.embed_image(_f4, w=175, h=120)

                if machine_col and machine_col in df_filt.columns and not df_filt[machine_col].dropna().empty:
                    _mc4 = df_filt[machine_col].value_counts()
                    _labels_mc = [f"{k}（{int(v)}件）" for k, v in _mc4.items()]
                    _f5, _a5 = _mplt.subplots(figsize=(7, 5))
                    wedges5, texts5, autotexts5 = _a5.pie(
                        list(_mc4.values),
                        labels=None,
                        autopct=lambda p: f"{p:.0f}%" if p >= 5 else "",
                        colors=["#FF5000","#060E9F","#8EB9C9","#FFCE00"],
                        startangle=90,
                        pctdistance=0.72,
                        wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
                    )
                    for _at5 in autotexts5:
                        _at5.set_fontsize(12)
                        _at5.set_fontweight("bold")
                    _a5.legend(
                        wedges5, _labels_mc,
                        loc="center left",
                        bbox_to_anchor=(1.0, 0.5),
                        fontsize=10,
                        frameon=False,
                    )
                    _a5.set_title(f"{period_label}　機台客訴佔比", fontsize=13, pad=12)
                    _f5.tight_layout()
                    pdf.embed_image(_f5, w=175, h=110)

                # ════════════════════════════════════════════
                # Page 5：趨勢 + 部門分析
                # ════════════════════════════════════════════
                pdf.add_page()
                pdf.section_title("客訴趨勢分析")

                _daily3 = df_filt.groupby(df_filt[date_col].dt.date).size().reset_index(name="件數")
                if len(_daily3) > 1:
                    _f6, _a6 = _mplt.subplots(figsize=(10, 4))
                    _a6.bar([str(d) for d in _daily3.iloc[:,0]], list(_daily3["件數"]),
                            color="#060E9F", edgecolor="white", linewidth=0.5)
                    _a6.set_title(f"{period_label}　每日件數趨勢", fontsize=13)
                    _a6.tick_params(axis="x", rotation=30, labelsize=8)
                    _a6.yaxis.set_major_locator(_MNL(integer=True))
                    _a6.set_ylabel("件數", fontsize=10)
                    _a6.grid(axis="y", alpha=0.3)
                    _f6.tight_layout()
                    pdf.embed_image(_f6, w=180, h=100)

                if dept_col and dept_col in df_filt.columns:
                    pdf.section_title("各部門件數分析")
                    _dp = df_filt[dept_col].replace("","未分配").value_counts()
                    rows_dp = [[i+1, str(d), int(v), f"{int(v)/n_cur*100:.0f}%"]
                                for i, (d, v) in enumerate(_dp.items())]
                    pdf.full_table(["排名","部門","件數","佔比"], rows_dp, [15,80,20,20])

                _pdf_bytes = bytes(pdf.output())
                st.session_state["_s4_pdf_bytes"] = _pdf_bytes
                st.session_state["_s4_pdf_label"] = period_label
                st.success(f"✅ PDF 已產生，共 {pdf.page_no()} 頁（{len(_pdf_bytes)//1024} KB）")
            except Exception as _e:
                import traceback
                st.error(f"PDF 產生失敗：{_e}")
                st.code(traceback.format_exc())

    if st.session_state.get("_s4_pdf_bytes"):
        _label = st.session_state.get("_s4_pdf_label", period_label)
        st.download_button(
            "⬇️ 下載完整分析 PDF（多頁）",
            data=st.session_state["_s4_pdf_bytes"],
            file_name=f"ECOCO_客訴分析_{_label.replace(' ','').replace('～','-').replace('/','-')}.pdf",
            mime="application/pdf",
            use_container_width=False,
            key="s4_dl_full_pdf",
        )

    # ── AI 口說報告 ────────────────────────────────────────────────
    st.markdown("---")
    st.markdown(f'<div class="s4-section">🎙️ AI 口說報告產生器</div>', unsafe_allow_html=True)
    rep_type = st.radio("報告類型", ["週會報告","月會報告","季報","年度報告"], horizontal=True, key="s4v3_rep")

    if st.button("🚀 產生 AI 口說報告", type="primary", key="s4v3_gen"):
        total_cur  = len(df_filt)
        total_prev = len(df_prev) if not df_prev.empty else None
        pct_chg    = pct_change(total_cur, total_prev) if total_prev else None

        type_summary = ""
        if type_col and type_col in df_filt.columns:
            _cvs = df_filt[type_col].value_counts()
            _pvs = df_prev[type_col].value_counts() if not df_prev.empty and type_col in df_prev.columns else pd.Series(dtype=int)
            for cat, cnt in _cvs.items():
                prev_cnt = int(_pvs.get(cat, 0))
                d = int(cnt) - prev_cnt
                pline = f"（較上期 {d:+d} 件，{pct_change(int(cnt),prev_cnt):+.1f}%）" if prev_cnt else ""
                type_summary += f"- {cat}：{int(cnt)} 件{pline}\n"

        city_summary = ""
        if city_col and city_col in df_filt.columns:
            _cc = df_filt[city_col].value_counts()
            _pc = df_prev[city_col].value_counts() if not df_prev.empty and city_col in df_prev.columns else pd.Series(dtype=int)
            for city, cnt in _cc.head(5).items():
                d = int(cnt) - int(_pc.get(city,0))
                city_summary += f"- {city}：{int(cnt)} 件（{d:+d}）\n"

        top3 = ""
        if detail_col and detail_col in df_filt.columns:
            for _, r in df_filt[detail_col].value_counts().head(3).reset_index().iterrows():
                top3 += f"- {r[detail_col]}：{r['count']} 件\n"

        _upper_cmp = (
            f"\n【上期對比】（{period_prev}，{total_prev} 件，總件數 {pct_chg:+.1f}%）"
            if pct_chg is not None else ""
        )
        prompt = (
            f"你是 ECOCO 宜可可循環經濟客服部的高級分析專員。\n"
            f"請根據以下數據，產出一份{rep_type}的「口說報告」，適合在會議中對長官簡報。\n\n"
            f"【語氣】：專業、條理清晰、帶有建議性，如現場口語報告。\n"
            f"【結構】：\n"
            f"1. 開場白（點出本期重點）\n"
            f"2. 總體趨勢概述（數字意義，非只念數字）\n"
            f"3. 前三大痛點深度解析（原因與影響）\n"
            f"4. 城市/區域分析亮點\n"
            f"5. 改善成效追蹤\n"
            f"6. 下階段行動建議\n\n"
            f"【本期數據】（{period_label}，共 {total_cur} 件）：\n"
            f"{type_summary or '（無問題類型資料）'}\n\n"
            f"【城市分布 TOP5】：\n"
            f"{city_summary or '（無城市資料）'}\n\n"
            f"【前三大問題細項】：\n"
            f"{top3 or '（無細項資料）'}\n"
            f"{_upper_cmp}\n\n"
            f"請以繁體中文撰寫，口語自然但不失專業，每段落 2-4 句。"
        )

        with st.spinner("AI 正在撰寫口說報告..."):
            from automation.llm import complete_text as _complete_text

            report_text = _complete_text(prompt, max_tokens=2000)
            if not report_text:
                try:
                    report_text = f"【口說報告】\n\n{generate_ai_summary(df_filt)}"
                except Exception as e:
                    report_text = f"⚠️ AI 暫時無法使用（{e}）\n\n數據摘要：\n\n{prompt}"

        st.text_area("📋 口說報告（可複製）", report_text, height=460, key="s4v3_report_out")

        dl_c1, dl_c2 = st.columns(2)
        dl_c1.download_button("⬇️ 下載口說報告（TXT）",
                              data=report_text.encode("utf-8"),
                              file_name=f"{period_label}_{rep_type}.txt",
                              mime="text/plain", key="s4v3_dl_txt",
                              use_container_width=True)
        try:
            _setup_cjk_font()
            import matplotlib.pyplot as _mplt
            _charts = {}
            if type_col and type_col in df_filt.columns:
                _tc2 = df_filt[type_col].value_counts()
                _f, _a = _mplt.subplots(figsize=(8,4))
                _colors = ["#060E9F","#FF5000","#FFCE00","#8EB9C9","#0076A9"]
                _a.pie(list(_tc2.values), labels=list(_tc2.index), autopct="%1.1f%%",
                       colors=_colors[:len(_tc2)], startangle=90)
                _a.set_title(f"{period_label} 客訴類別分佈")
                _b = io.BytesIO(); _f.savefig(_b, format="png", dpi=150, bbox_inches="tight")
                _mplt.close(_f); _charts["客訴類別分佈.png"] = _b.getvalue()
            if city_col and city_col in df_filt.columns:
                _cc2 = df_filt[city_col].value_counts().head(10)
                _f2, _a2 = _mplt.subplots(figsize=(8,4))
                _a2.bar(list(_cc2.index), list(_cc2.values), color="#FF5000")
                _a2.set_title(f"{period_label} 城市件數排行")
                _a2.yaxis.set_major_locator(_mplt.MaxNLocator(integer=True))
                _a2.tick_params(axis="x", rotation=15)
                _b2 = io.BytesIO(); _f2.savefig(_b2, format="png", dpi=150, bbox_inches="tight")
                _mplt.close(_f2); _charts["城市排行.png"] = _b2.getvalue()
            if detail_col and detail_col in df_filt.columns:
                _dc2 = df_filt[detail_col].value_counts().head(8)
                _f3, _a3 = _mplt.subplots(figsize=(8,4))
                _a3.barh(list(_dc2.index)[::-1], list(_dc2.values)[::-1], color="#060E9F")
                _a3.set_title(f"{period_label} TOP 8 問題細項")
                _b3 = io.BytesIO(); _f3.savefig(_b3, format="png", dpi=150, bbox_inches="tight")
                _mplt.close(_f3); _charts["問題細項排行.png"] = _b3.getvalue()
            _zbuf = io.BytesIO()
            with zipfile.ZipFile(_zbuf, "w", zipfile.ZIP_DEFLATED) as _zf:
                for _fn, _fb in _charts.items():
                    _zi = zipfile.ZipInfo(_fn); _zi.flag_bits |= 0x800
                    _zi.compress_type = zipfile.ZIP_DEFLATED; _zf.writestr(_zi, _fb)
                _zr = zipfile.ZipInfo(f"{period_label}_{rep_type}.txt")
                _zr.flag_bits |= 0x800; _zf.writestr(_zr, report_text.encode("utf-8"))
            dl_c2.download_button("⬇️ 下載圖表+報告（ZIP）",
                                  data=_zbuf.getvalue(),
                                  file_name=f"{period_label}_趨勢分析.zip",
                                  mime="application/zip", key="s4v3_dl_zip",
                                  use_container_width=True)
        except Exception as _ze:
            dl_c2.warning(f"ZIP 產生失敗：{_ze}")



# (session key, 顯示標籤, 前面是否加分隔線)
# 標籤用一般半形空格，不要用全形空格：26px 下全形空格就吃掉一個字的寬度，
# 「客訴趨勢分析儀表板」會被擠到折行。
NAV_ITEMS = [
    ("首頁", "🏠 首頁", False),
    ("上傳檔案區（分析區）", "⬆️ 上傳檔案區 (分析區)", False),
    ("圖表與 AI 分析", "🔎 圖表與 AI 分析", False),
    ("歷史紀錄", "🕘 歷史紀錄", False),
    ("趨勢分析", "📈 客訴趨勢分析儀表板", True),
]


def render_sidebar_nav() -> str:
    """左側導覽；主標題由 CSS 放在頂部橫條，不佔側邊欄空間。"""
    if "menu" not in st.session_state:
        st.session_state["menu"] = "首頁"
    if st.session_state["menu"] == "功能列表區":   # 舊 session 的值
        st.session_state["menu"] = "首頁"

    with st.sidebar:
        for key, label, divider in NAV_ITEMS:
            if divider:
                st.markdown("<hr>", unsafe_allow_html=True)
            active = st.session_state["menu"] == key
            if st.button(label, key=f"nav_{key}", use_container_width=True,
                         type="primary" if active else "secondary"):
                st.session_state["menu"] = key
                st.rerun()
    return st.session_state["menu"]


# 首頁功能卡：(圖示 svg path, 標題, 條列內容)
HOME_CARDS = [
    (
        "M14 4H8a2 2 0 0 0-2 2v20a2 2 0 0 0 2 2h9M14 4l6 6M14 4v6h6"
        "M20 10v6M24 28l4-4-4-4M28 24h-9",
        "多樣化檔案上傳與編輯",
        [
            "上傳 excel/csv/pdf，分析並標記問題<b>【問題類型、問題細項】</b>；",
            "支援下拉選填、編輯、篩選；",
            "批次勾選編輯/刪除；",
            "下載 Excel、上傳 Google Sheet。",
        ],
    ),
    (
        "M16 4a12 12 0 1 0 12 12h-12V4Z M20 4a12 12 0 0 1 8 8h-8V4Z",
        "視覺化圖表與部門分析",
        [
            "將分析結果圖表化；",
            "顯示各類型件數與百分比；",
            "並標示歸屬部門；",
            "可預覽與下載 AI 重點分析。",
        ],
    ),
    (
        "M16 6a10 10 0 1 1-9.8 12M6 6v6h6M16 11v5l4 3",
        "歷史紀錄管理",
        [
            "歷史分析紀錄管理 (最新置頂)，",
            "可預覽與下載。",
        ],
    ),
    (
        "M4 24l7-8 5 4 7-9M23 11h5v5M4 28h24",
        "多維度趨勢分析儀表板",
        [
            "週/月/季/年度趨勢分析；",
            "從歷史紀錄聚合數據，趨勢對比；",
            "AI 口說報告產生器。",
        ],
    ),
]


def render_home_cards() -> None:
    cards = []
    for path, title, bullets in HOME_CARDS:
        items = "".join(f"<li>{b}</li>" for b in bullets)
        cards.append(
            "<div class='home-card'>"
            "<div class='home-card-icon'>"
            "<svg viewBox='0 0 32 32' fill='none' stroke='currentColor' "
            "stroke-width='1.8' stroke-linecap='round' stroke-linejoin='round'>"
            f"<path d='{path}'/></svg></div>"
            f"<div><div class='home-card-title'>{title}</div><ul>{items}</ul></div>"
            "</div>"
        )
    st.markdown(f"<div class='home-grid'>{''.join(cards)}</div>", unsafe_allow_html=True)


def empty_state(message: str = "沒有資料紀錄") -> None:
    """查無資料時的中性提示；不要用紅色錯誤框嚇人。"""
    st.markdown(
        "<div style='background:#f6f8fb;border:1px solid #dfe5ee;border-radius:10px;"
        "padding:18px 22px;color:#5a6472;text-align:center;margin:10px 0;'>"
        f"{message}</div>",
        unsafe_allow_html=True,
    )


def page_header(title: str, subtitle: str = "") -> None:
    """頁首：淺底大標題 + 小字副標，不加色塊。"""
    sub = f"<div class='page-header-sub'>{subtitle}</div>" if subtitle else ""
    st.markdown(
        f"<div class='page-header'><div class='page-header-title'>{title}</div>{sub}</div>",
        unsafe_allow_html=True,
    )


def main():
    apply_brand_theme()
    menu = render_sidebar_nav()

    if menu == "首頁":
        page_header("首頁", "客訴分析平台功能總覽")
        render_home_cards()
    elif menu == "上傳檔案區（分析區）":
        section_1()
    elif menu == "圖表與 AI 分析":
        section_2()
    elif menu == "趨勢分析":
        section_4()
    else:
        section_3()

    # 版權頁尾已移除；只留回到頁首的按鈕。
    st.markdown(
        """
        <style>
            .scroll-top-btn {
                position: fixed;
                left: 18px;
                bottom: 18px;
                z-index: 1000;
                border: 1px solid #8EB9C9;
                background: #FFFFFF;
                color: #060E9F;
                border-radius: 999px;
                padding: 10px 18px;
                font-size: 15px;
                cursor: pointer;
                box-shadow: 0 2px 8px rgba(0,0,0,.12);
            }
        </style>
        <button class="scroll-top-btn" onclick="window.parent.scrollTo({top:0, behavior:'smooth'});">⌃&nbsp; 置頂</button>
        """,
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    main()
